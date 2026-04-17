# -*- coding: utf-8 -*-
"""웹 리서치 어시스턴트 — 터미널 버전 v2"""

import os, sys, json, re, time
from datetime import datetime

VERSION = "2.5.2"
_GITHUB_RAW = "https://raw.githubusercontent.com/caifyhelp-cmyk/web-researcher/master"

def _check_update():
    """GitHub에서 최신 버전 확인 후 자동 업데이트"""
    try:
        import urllib.request
        url = f"{_GITHUB_RAW}/version.txt"
        with urllib.request.urlopen(url, timeout=4) as r:
            latest = r.read().decode().strip()
        if latest == VERSION:
            return
        print(f"  업데이트 발견: {VERSION} → {latest}  다운로드 중...")
        here = os.path.dirname(os.path.abspath(__file__))
        for gh_name, local_name in [("app_local.py",          os.path.basename(__file__)),
                                     ("web_researcher.py",    "web_researcher.py"),
                                     ("orchestrator.py",      "orchestrator.py"),
                                     ("feedback_collector.py","feedback_collector.py"),
                                     ("maestro.py",           "maestro.py"),
                                     ("pattern_collector.py", "pattern_collector.py"),
                                     ("tools_kb.json",        "tools_kb.json")]:
            try:
                tmp, _ = urllib.request.urlretrieve(f"{_GITHUB_RAW}/{gh_name}")
                with open(tmp, "rb") as src, \
                     open(os.path.join(here, local_name), "wb") as dst:
                    dst.write(src.read())
            except Exception:
                pass
        print("  업데이트 완료. 자동 재시작합니다...\n")
        import subprocess, time
        time.sleep(1)
        subprocess.Popen([sys.executable, os.path.abspath(__file__)] + sys.argv[1:],
                         creationflags=subprocess.CREATE_NEW_CONSOLE if sys.platform == "win32" else 0)
        sys.exit(0)
    except Exception:
        pass  # 오프라인이거나 실패 시 그냥 실행

_check_update()

# API 키 자동 주입
try:
    import _local_keys
except ImportError:
    pass

# ── Rich UI ──────────────────────────────────────────────────────
from rich.console import Console
from rich.panel import Panel
from rich.progress import Progress, SpinnerColumn, TextColumn, TimeElapsedColumn
from rich.table import Table
from rich.rule import Rule
from rich.prompt import Prompt, Confirm

# ── PDF / PPT ────────────────────────────────────────────────────
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    _PDF = True
except Exception:
    _PDF = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPT = True
except Exception:
    _PPT = False

import openpyxl
from openai import OpenAI
from anthropic import Anthropic

# ── API 클라이언트 ────────────────────────────────────────────────
OPENAI_KEY    = os.getenv("OPENAI_API_KEY", "")
ANTHROPIC_KEY = os.getenv("ANTHROPIC_API_KEY", "")
DEEPSEEK_KEY  = os.getenv("DEEPSEEK_API_KEY", "")
GROK_KEY      = os.getenv("GROK_API_KEY", "")

oai      = OpenAI(api_key=OPENAI_KEY)                                         if OPENAI_KEY    else None
claude_c = Anthropic(api_key=ANTHROPIC_KEY)                                   if ANTHROPIC_KEY else None
deepseek = OpenAI(api_key=DEEPSEEK_KEY, base_url="https://api.deepseek.com") if DEEPSEEK_KEY  else None
grok_ai  = OpenAI(api_key=GROK_KEY,     base_url="https://api.x.ai/v1")      if GROK_KEY      else None

console  = Console(highlight=False)

# ── 오케스트레이터 ────────────────────────────────────────────────
try:
    import orchestrator as orch
    _ORCH_ENABLED = True
except Exception:
    _ORCH_ENABLED = False

# ── 피드백 수집기 ─────────────────────────────────────────────────
try:
    import feedback_collector as fc
    _FC_ENABLED = True
except Exception:
    _FC_ENABLED = False

_session_models: dict = {}  # category -> model_name (이번 세션 추적용)


# ── 백그라운드 자기평가 ───────────────────────────────────────────
def _start_background_evaluation():
    """앱 시작 시 7일 이상 된 카테고리를 백그라운드에서 재평가"""
    if not _ORCH_ENABLED:
        return

    stale = orch.get_stale_categories(days=7)
    if not stale:
        return

    import threading

    def _run():
        # Meta caller (DeepSeek → GPT-4o fallback)
        def meta_caller(prompt: str) -> str:
            if deepseek:
                try:
                    r = deepseek.chat.completions.create(
                        model="deepseek-chat",
                        messages=[{"role": "user", "content": prompt}],
                        max_tokens=300
                    )
                    return r.choices[0].message.content.strip()
                except Exception:
                    pass
            return call_gpt(prompt, model="gpt-4o", max_tokens=300)

        # Self-callers
        self_callers = {}
        if oai:
            self_callers["gpt-4o"]      = lambda p: call_gpt(p, model="gpt-4o",      max_tokens=50)
            self_callers["gpt-4o-mini"] = lambda p: call_gpt(p, model="gpt-4o-mini", max_tokens=50)
        if deepseek:
            self_callers["deepseek"] = lambda p: call_deepseek(p)
        if claude_c:
            self_callers["claude"] = lambda p: call_claude(p)

        if not self_callers:
            return

        for cat in stale:
            try:
                winner = orch.run_self_evaluation(cat, meta_caller, self_callers)
                orch.reset_streak(cat)
            except Exception:
                pass

    t = threading.Thread(target=_run, daemon=True)
    t.start()

# 저장 폴더 — Desktop 없으면 Documents, 그것도 없으면 홈
def _get_save_dir():
    home = os.path.expanduser("~")
    for candidate in [
        os.path.join(home, "Desktop", "리서치결과"),
        os.path.join(home, "바탕 화면", "리서치결과"),
        os.path.join(home, "Documents", "리서치결과"),
        os.path.join(home, "리서치결과"),
    ]:
        try:
            os.makedirs(candidate, exist_ok=True)
            return candidate
        except Exception:
            continue
    return home

SAVE_DIR = _get_save_dir()

# ── 뇌 에이전트 연동 ─────────────────────────────────────────────
_BRAIN_AGENT_URL = "https://brain-agent-v9wl.onrender.com/api/research"
_BRAIN_AGENT = None

def _get_brain_agent():
    """로컬 thinking_agent import 시도 (개발자 PC 전용)"""
    global _BRAIN_AGENT
    if _BRAIN_AGENT is None:
        try:
            brain_path = os.path.join(os.path.expanduser("~"), "thinking_agent")
            if brain_path not in sys.path:
                sys.path.insert(0, brain_path)
            from agent import analyze as _brain_analyze
            _BRAIN_AGENT = _brain_analyze
        except Exception:
            _BRAIN_AGENT = False
    return _BRAIN_AGENT if _BRAIN_AGENT else None

def call_brain_agent(situation: str) -> str:
    """뇌 에이전트 호출 — HTTP API 우선, 로컬 폴백, 실패 시 빈 문자열"""
    # 1순위: HTTP API (모든 PC에서 작동)
    try:
        resp = requests.post(
            _BRAIN_AGENT_URL,
            json={"situation": situation},
            timeout=15
        )
        if resp.status_code == 200:
            data = resp.json()
            if data.get("ok"):
                parts = []
                if data.get("judgment"): parts.append(f"판단: {data['judgment']}")
                if data.get("action"):   parts.append(f"액션: {data['action']}")
                if data.get("reason"):   parts.append(f"근거: {data['reason']}")
                return "\n".join(parts)
    except Exception:
        pass

    # 2순위: 로컬 import (개발자 PC)
    fn = _get_brain_agent()
    if fn:
        try:
            result = fn(situation)
            parts = []
            if result.get("judgment"): parts.append(f"판단: {result['judgment']}")
            if result.get("action"):   parts.append(f"액션: {result['action']}")
            if result.get("reason"):   parts.append(f"근거: {result['reason']}")
            return "\n".join(parts)
        except Exception:
            pass

    return ""

# ═══════════════════════════════════════════════════════
#  LLM 호출
# ═══════════════════════════════════════════════════════
def call_gpt(prompt: str, system: str = "", model: str = "gpt-4o-mini",
             temperature: float = 1.0, max_tokens: int = 2000) -> str:
    if not oai:
        return "[OpenAI 키 없음]"
    msgs = ([{"role": "system", "content": system}] if system else []) + \
           [{"role": "user", "content": prompt}]
    try:
        r = oai.chat.completions.create(model=model, messages=msgs,
                                        temperature=temperature, max_tokens=max_tokens)
        return r.choices[0].message.content.strip()
    except Exception as e:
        return f"[GPT 오류: {e}]"


def call_model(category: str, prompt: str, system: str = "",
               temperature: float = 1.0, max_tokens: int = 2000) -> str:
    """오케스트레이터가 선택한 최적 모델로 동적 라우팅."""
    if _ORCH_ENABLED:
        # streak >= 3이면 백그라운드 재평가 트리거
        if orch.check_reevaluation_needed(category):
            _start_background_evaluation()
        model = orch.get_best_model(category)
    else:
        model = {"query_generation": "deepseek", "url_filtering": "gpt-4o",
                 "data_extraction": "gpt-4o", "market_analysis": "gpt-4o",
                 "strategy_insight": "claude"}.get(category, "gpt-4o-mini")

    _session_models[category] = model

    if model == "claude":
        return call_claude(prompt, system)
    elif model == "gpt-4o":
        return call_gpt(prompt, system, model="gpt-4o",
                        temperature=temperature, max_tokens=max_tokens)
    elif model == "deepseek":
        return call_deepseek(prompt)
    else:
        return call_gpt(prompt, system, model="gpt-4o-mini",
                        temperature=temperature, max_tokens=max_tokens)

def call_claude(prompt: str, system: str = "") -> str:
    if not claude_c:
        return call_gpt(prompt, system)
    try:
        r = claude_c.messages.create(
            model="claude-opus-4-6", max_tokens=2000,
            system=system or "당신은 전문 리서치 분석가입니다.",
            messages=[{"role": "user", "content": prompt}]
        )
        return r.content[0].text.strip()
    except Exception as e:
        return f"[Claude 오류: {e}]"

def call_deepseek(prompt: str) -> str:
    if not deepseek:
        return call_gpt(prompt)
    try:
        r = deepseek.chat.completions.create(
            model="deepseek-reasoner",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=2000
        )
        return r.choices[0].message.content.strip()
    except Exception:
        return call_gpt(prompt)

# ═══════════════════════════════════════════════════════
#  리서치 플랜 (DeepSeek) — 조사 유형 + 필터 모드 자동 판단
# ═══════════════════════════════════════════════════════
def make_plan(topic: str, depth: str) -> dict:
    prompt = f"""리서치 주제: {topic}
조사 깊이: {depth}

━━ 1단계: 조사 유형 판단 ━━
- "competitor"  : 경쟁사·업체·서비스 비교 조사
- "news"        : 뉴스·트렌드·이슈·최신 동향 조사
- "institution" : 특정 기관·단체·기업 심층 분석
- "general"     : 일반 정보 수집

━━ 2단계: 필터 모드 판단 ━━
- "strict" : 서비스 제공 업체·기관만 (뉴스·블로그·외국사이트 차단)
- "medium" : 관련 사이트 폭넓게 (블로그 허용, 뉴스 제한)
- "loose"  : 뉴스·블로그·해외사이트 모두 허용

기본 규칙: competitor→strict / news→loose / institution→medium / general→medium

━━ 3단계: 검색어 생성 ━━
핵심 원칙: 고객이 요청한 주제의 핵심 단어를 반드시 포함할 것. 임의로 주제를 확장하거나 다른 영역으로 벗어나지 말 것.
- 모든 검색어에 리서치 주제의 핵심 키워드가 포함되어야 함
- 각 검색어는 같은 주제를 다른 각도로 표현 (의미가 완전히 달라지면 안 됨)
- 수식어만 바꾸기: "비교", "추천", "업체 목록", "후기", "가격" 등
- news 유형이면 "최신", "2025", "동향", "트렌드" 등 시의성 키워드 추가
- 주제와 관련 없는 단어나 상위 개념으로 확장 금지

━━ 4단계: needs 항목 ━━
- 주제에 맞는 실질적 추출 항목 5~8개

JSON만 답하세요:
{{
  "research_type": "competitor|news|institution|general",
  "filter_mode": "strict|medium|loose",
  "summary": "주제 한줄 요약",
  "keywords": ["핵심키워드1", "핵심키워드2"],
  "queries": ["검색어1", "검색어2", "검색어3", "검색어4", "검색어5"],
  "needs": ["항목1", "항목2", "항목3", "항목4", "항목5"],
  "focus_points": ["핵심1", "핵심2", "핵심3"],
  "analysis_angle": "분석 방향"
}}"""
    raw = call_deepseek(prompt)
    try:
        m = re.search(r'\{.*\}', raw, re.DOTALL)
        if m:
            plan = json.loads(m.group())
            rtype = plan.get("research_type", "general")
            fmode = plan.get("filter_mode", "medium")
            type_label = {"competitor":"경쟁사 조사","news":"뉴스·트렌드","institution":"기관 분석","general":"일반 조사"}.get(rtype, rtype)
            filter_label = {"strict":"엄격 (업체만)","medium":"보통","loose":"느슨 (뉴스·블로그 허용)"}.get(fmode, fmode)
            console.print(f"[dim]유형: {type_label}  |  필터: {filter_label}[/dim]")
            console.print(f"[dim]생성된 검색어:[/dim]")
            for i, q in enumerate(plan.get("queries", []), 1):
                console.print(f"[dim]  {i}. {q}[/dim]")
            return plan
    except Exception:
        pass
    return {
        "research_type": "general", "filter_mode": "medium",
        "summary": topic, "keywords": [topic],
        "queries": [topic, f"{topic} 비교", f"{topic} 업체", f"{topic} 현황", f"{topic} 사례"],
        "needs": ["업체명", "주요서비스", "가격정책", "고객후기", "연락처"],
        "focus_points": ["주요 현황", "서비스 분석", "시장 트렌드"],
        "analysis_angle": "종합 전략 분석"
    }

# ═══════════════════════════════════════════════════════
#  URL 관련성 필터 — filter_mode에 따라 기준 조정
# ═══════════════════════════════════════════════════════
def _filter_relevant(candidates: list, topic: str, keywords: list, filter_mode: str = "medium") -> list:
    if not candidates or not oai:
        return candidates
    items = "\n".join(f"{i+1}. [{c.get('title','')}] {c.get('url','')}"
                      for i, c in enumerate(candidates))
    kw = ", ".join(keywords) if keywords else topic

    if filter_mode == "strict":
        criteria = (
            "제거 기준:\n"
            "- 뉴스·언론·미디어 사이트\n"
            "- 블로그·개인 사이트\n"
            "- 정부 부처·공공데이터 포털\n"
            "- 대학교·학과 사이트\n"
            "- 한국 서비스가 아닌 외국 사이트\n"
            "- SNS·쇼핑몰·부동산 등 무관 업종\n"
            "포함 기준: 한국에서 해당 서비스를 직접 제공하는 업체·기관·협회"
        )
    elif filter_mode == "loose":
        criteria = (
            "제거 기준:\n"
            "- 주제와 완전히 무관한 사이트 (쇼핑몰, 부동산, 식품 등)\n"
            "- SNS 플랫폼 (인스타, 유튜브 등)\n"
            "포함 기준: 뉴스·블로그·해외사이트도 주제와 관련 있으면 포함"
        )
    else:  # medium
        criteria = (
            "제거 기준:\n"
            "- 뉴스·언론 사이트\n"
            "- 정부 부처·공공데이터 포털\n"
            "- 대학교·학과 사이트\n"
            "- 한국 서비스가 아닌 외국 사이트\n"
            "- SNS·쇼핑몰 등 완전히 무관 업종\n"
            "포함 기준: 해당 주제와 직접 관련된 기관·업체·블로그·리뷰"
        )

    resp = call_model(
        "url_filtering",
        f"주제: {topic}\n핵심 키워드: {kw}\n\n{criteria}\n\n{items}\n\n"
        f"제거할 번호만 쉼표로 (없으면 '없음'):",
        system="URL 관련성 판단기. 제거할 번호만 답하세요.",
    )
    try:
        if "없음" in resp or not re.search(r'\d', resp):
            return candidates
        remove_nums = {int(x.strip()) - 1 for x in re.findall(r'\d+', resp)}
        filtered = [c for i, c in enumerate(candidates) if i not in remove_nums]
        return filtered if filtered else candidates
    except Exception:
        return candidates

# ═══════════════════════════════════════════════════════
#  스크래퍼 (requests + Selenium 폴백)
# ═══════════════════════════════════════════════════════
import requests
from bs4 import BeautifulSoup

_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/124.0 Safari/537.36",
    "Accept-Language": "ko-KR,ko;q=0.9",
    "Referer": "https://www.naver.com",
}

def _safe_decode(resp) -> str:
    candidates = []
    try:
        from charset_normalizer import from_bytes
        detected = from_bytes(resp.content).best()
        if detected:
            candidates.append(str(detected.encoding))
    except Exception:
        pass
    try:
        if resp.apparent_encoding:
            candidates.append(resp.apparent_encoding)
    except Exception:
        pass
    candidates += ["utf-8", "cp949", "euc-kr"]
    for enc in dict.fromkeys(candidates):
        try:
            return resp.content.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return resp.content.decode("utf-8", errors="replace")

def _parse_html(html: str) -> dict:
    soup = BeautifulSoup(html, "html.parser")
    page_title = soup.title.string.strip() if soup.title and soup.title.string else ""
    headings = " | ".join(
        h.get_text(strip=True)
        for h in soup.find_all(["h1","h2","h3"])
        if h.get_text(strip=True)
    )[:500]
    for tag in soup(["script","style","nav","footer","header","aside"]):
        tag.decompose()
    text = " ".join(soup.get_text(separator=" ").split())[:4000]
    return {"text": text, "headings": headings, "page_title": page_title, "error": ""}

_selenium_driver = None

def _get_driver():
    global _selenium_driver
    if _selenium_driver is None:
        try:
            from web_researcher import make_driver
            console.print("[dim]Selenium 드라이버 초기화 중...[/dim]")
            _selenium_driver = make_driver()
        except Exception:
            pass
    return _selenium_driver

def _fetch_page(url: str, timeout: int = 8) -> dict:
    empty = {"text": "", "headings": "", "page_title": "", "error": ""}
    last_error = ""

    for attempt in range(3):
        try:
            resp = requests.get(url, headers=_HEADERS, timeout=timeout, allow_redirects=True)
            resp.raise_for_status()
            ct = resp.headers.get("Content-Type", "").lower()
            if ct and "html" not in ct and "plain" not in ct:
                return {**empty, "error": "non-html content-type"}
            result = _parse_html(_safe_decode(resp))
            if len(result["text"]) >= 150:
                return result
            last_error = "JS rendering detected"
            break
        except Exception as e:
            last_error = str(e)[:80]
            if attempt < 2:
                time.sleep(1)

    try:
        driver = _get_driver()
        if driver:
            from web_researcher import scrape_page
            page = scrape_page(driver, url)
            if page.get("full_text") and len(page["full_text"]) > 150:
                return {
                    "text":       page["full_text"][:4000],
                    "headings":   page.get("headings", ""),
                    "page_title": page.get("page_title", ""),
                    "error":      "",
                }
    except Exception as e:
        last_error = str(e)[:80]

    return {**empty, "error": last_error}

def _fetch_text(url: str, timeout: int = 8) -> str:
    return _fetch_page(url, timeout)["text"]

# ═══════════════════════════════════════════════════════
#  검색 + 스크랩
# ═══════════════════════════════════════════════════════
def _collect_urls(queries: list, plan: dict, seen_domains: set, prog=None) -> list:
    try:
        from web_researcher import search_naver, search_duckduckgo, HARD_EXCLUDE, BLOG_PLATFORM_EXCLUDE
    except ImportError:
        console.print("[red]web_researcher 모듈 없음[/red]")
        return []

    from urllib.parse import urlparse
    results  = []
    keywords = plan.get("keywords", [])
    fmode    = plan.get("filter_mode", "medium")

    # filter_mode에 따라 차단 목록 조정
    if fmode == "strict":
        blocked = set(HARD_EXCLUDE + BLOG_PLATFORM_EXCLUDE)
    elif fmode == "loose":
        blocked = set(HARD_EXCLUDE)  # SNS·포털만 차단, 뉴스·블로그 허용
    else:  # medium
        blocked = set(HARD_EXCLUDE)  # BLOG_PLATFORM은 GPT 필터에 맡김

    # 한글 비율 기준: loose면 완화 (5%), 나머지 15%
    korean_threshold = 0.05 if fmode == "loose" else 0.15

    def _is_blocked(domain: str) -> bool:
        bare = domain.replace("www.", "").replace("m.", "")
        return any(bare == b or bare.endswith("." + b) for b in blocked)

    for q in queries:
        task_id = prog.add_task(f"[cyan]검색: {q[:45]}", total=None) if prog else None
        try:
            naver = search_naver(q, max_results=10)
            candidates = naver
            if len(naver) < 3:
                ddg = search_duckduckgo(q, max_results=8)
                naver_urls = {n["url"] for n in naver}
                candidates = naver + [d for d in ddg if d["url"] not in naver_urls]

            candidates = [c for c in candidates if not _is_blocked(c.get("domain", ""))]
            candidates = _filter_relevant(candidates, q, keywords, fmode)

            for item in candidates[:5]:
                url   = item.get("url", "")
                title = item.get("title", url)
                if not url:
                    continue
                domain = urlparse(url).netloc
                if domain in seen_domains:
                    continue
                seen_domains.add(domain)
                if prog and task_id is not None:
                    prog.update(task_id, description=f"[yellow]수집: {title[:40]}")
                page = _fetch_page(url)
                text = page["text"]
                if not text or len(text) < 150:
                    continue
                korean_ratio = len(re.findall(r'[가-힣]', text)) / len(text)
                if korean_ratio < korean_threshold:
                    continue
                results.append({
                    "query":      q,
                    "url":        url,
                    "title":      title,
                    "domain":     domain,
                    "content":    text[:3000],
                    "page_title": page["page_title"],
                    "headings":   page["headings"],
                    "error":      page["error"],
                })
        except Exception:
            pass
        if prog and task_id is not None:
            prog.remove_task(task_id)

    return results


def run_research(plan: dict) -> list:
    from urllib.parse import urlparse
    seen_domains = set()
    queries = plan.get("queries", [])
    depth_n = {"빠른 조사": 3, "일반 조사": 5, "심층 조사": 8}
    max_q   = depth_n.get(plan.get("_depth", "일반 조사"), 5)
    target  = {"빠른 조사": 5, "일반 조사": 10, "심층 조사": 18}.get(plan.get("_depth", "일반 조사"), 10)

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        TimeElapsedColumn(),
        console=console, transient=True
    ) as prog:
        results = _collect_urls(queries[:max_q], plan, seen_domains, prog)

        retry = 0
        while len(results) < int(target * 0.7) and retry < 2:
            retry += 1
            collected_domains = [r["domain"] for r in results]
            extra_prompt = (
                f'조사 목적: "{plan.get("summary", "")}"\n'
                f'기존 수집 도메인: {collected_domains[:10]}\n'
                f'아직 {target - len(results)}개가 더 필요합니다.\n\n'
                f'규칙:\n'
                f'- 실제 고객이 네이버에 검색할 자연스러운 검색어\n'
                f'- 기존 도메인과 겹치지 않는 새로운 각도의 쿼리 3개\n\n'
                f'JSON: {{"queries": ["q1", "q2", "q3"]}}'
            )
            try:
                raw = call_gpt(extra_prompt, model="gpt-4o-mini")
                m = re.search(r'\{[\s\S]+\}', raw)
                if m:
                    extra_queries = json.loads(m.group()).get("queries", [])
                    extra = _collect_urls(extra_queries, plan, seen_domains, prog)
                    results.extend(extra)
            except Exception:
                break

    return results

# ═══════════════════════════════════════════════════════
#  URL별 GPT 구조화 추출
# ═══════════════════════════════════════════════════════
def _analyze_url(r: dict, topic: str, needs: list) -> dict:
    if not oai:
        return {"한줄요약": "", **{n: "" for n in needs}}
    needs_str = "\n".join(f'  "{n}": ""' for n in needs)
    prompt = f"""웹 리서치 전문가로서 아래 페이지를 분석해주세요.

[조사 주제]: {topic}
[URL]: {r.get('url','')}
[제목]: {r.get('page_title', r.get('title',''))}
[헤딩]: {r.get('headings','')}
[본문]: {r.get('content','')}

없으면 "확인 불가". 각 항목 2~5문장.
JSON만:
{{
  "한줄요약": "",
{needs_str}
}}"""
    try:
        raw = call_model("data_extraction", prompt, temperature=0.3, max_tokens=1500)
        m = re.search(r"\{[\s\S]+\}", raw)
        if m:
            return json.loads(m.group())
    except Exception as e:
        console.print(f"[red]GPT 추출 오류: {e}[/red]")
    return {"한줄요약": "분석 실패", **{n: "" for n in needs}}

# ═══════════════════════════════════════════════════════
#  AI 분석
# ═══════════════════════════════════════════════════════
def analyze(topic: str, plan: dict, results: list) -> dict:
    if not results:
        return {"gpt_analysis": "수집된 데이터가 없습니다.",
                "claude_insights": "", "brain_insights": "",
                "source_count": 0, "sources": [], "per_url": []}

    needs = plan.get("needs", ["업체명", "주요서비스", "가격정책", "고객후기", "연락처"])

    # URL별 구조화 추출
    per_url = []
    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console, transient=True
    ) as prog:
        t = prog.add_task("[cyan]URL별 GPT 분석 중...", total=len(results))
        for r in results:
            gpt = _analyze_url(r, topic, needs)
            values = [v for k, v in gpt.items() if k != "한줄요약"]
            unconfirmed = sum(1 for v in values if str(v).strip() in ("확인 불가", "", "분석 실패"))
            if gpt.get("한줄요약") in ("분석 실패", "확인 불가") or unconfirmed >= len(values) * 0.5:
                prog.advance(t)
                continue
            per_url.append({**r, "gpt": gpt})
            prog.advance(t)

    ctx = "\n\n".join(
        f"[{r['title']}]\n출처: {r['url']}\n{r['content'][:1000]}"
        for r in results[:8]
    )

    gpt_out = call_model(
        "market_analysis",
        f"리서치 주제: {topic}\n\n수집 데이터:\n{ctx}\n\n"
        f"핵심 분석, 주요 플레이어, 시장 현황, 시사점을 한국어로 작성하세요.",
        system="당신은 시니어 마케팅 리서치 애널리스트입니다. 핵심만 간결하게.",
    )

    # 뇌 에이전트 보호: strategy/competitor/marketing 관련일 때만 호출
    # 오케스트레이터 학습 데이터는 절대 뇌 에이전트로 흘러가지 않음
    _BRAIN_MARKETING_KW = ["마케팅", "광고", "브랜드", "전략", "포지셔닝", "타겟", "캠페인", "홍보"]
    rtype = plan.get("research_type", "general")
    _should_brain = (
        rtype in ("competitor", "general") or
        any(kw in topic for kw in _BRAIN_MARKETING_KW)
    )

    brain_out = ""
    if _should_brain:
        with console.status("[dim]뇌 에이전트 분석 중...[/dim]"):
            brain_situation = (
                f"리서치 주제: {topic}\n"
                f"조사 유형: {rtype}\n"
                f"수집된 주요 정보 요약:\n{gpt_out[:800]}"
            )
            brain_out = call_brain_agent(brain_situation)

    # strategy_insight 동적 라우팅
    if brain_out:
        claude_out = call_model(
            "strategy_insight",
            f"주제: {topic}\n\nGPT 분석:\n{gpt_out}\n\n뇌 에이전트 판단:\n{brain_out}\n\n"
            f"위 내용을 종합해 실행 가능한 전략 액션 5개를 번호 목록으로 제시하세요.",
            system="마케팅 전략 전문가로서 실용적인 인사이트를 제공하세요."
        )
    else:
        claude_out = call_model(
            "strategy_insight",
            f"주제: {topic}\n\nGPT 분석 결과:\n{gpt_out}\n\n"
            f"전략적 시사점과 실행 가능한 액션 아이템 5개를 번호 목록으로 제시하세요.",
            system="마케팅 전략 전문가로서 실용적인 인사이트를 제공하세요."
        )

    return {
        "gpt_analysis":    gpt_out,
        "claude_insights": claude_out,
        "brain_insights":  brain_out,
        "source_count":    len(results),
        "sources":         [{"title": r["title"], "url": r["url"]} for r in results[:10]],
        "per_url":         per_url,
        "needs":           needs,
    }

# ═══════════════════════════════════════════════════════
#  결과 출력
# ═══════════════════════════════════════════════════════
def display_results(topic: str, plan: dict, analysis: dict):
    console.print()
    console.print(Rule(f"[bold cyan]리서치 결과: {topic}[/bold cyan]", style="bright_blue"))

    rtype = plan.get("research_type", "general")
    fmode = plan.get("filter_mode", "medium")
    type_label   = {"competitor":"경쟁사 조사","news":"뉴스·트렌드","institution":"기관 분석","general":"일반 조사"}.get(rtype, rtype)
    filter_label = {"strict":"엄격","medium":"보통","loose":"느슨"}.get(fmode, fmode)

    console.print(Panel(
        f"[bold]요약[/bold]  {plan.get('summary', topic)}\n"
        f"[bold]유형[/bold]  {type_label}  |  [bold]필터[/bold]  {filter_label}\n"
        f"[bold]수집[/bold]  {analysis['source_count']}개  →  분석 완료 {len(analysis['per_url'])}개\n"
        f"[bold]방향[/bold]  {plan.get('analysis_angle','')}",
        title="[bold]개요[/bold]", border_style="blue", padding=(0,2)
    ))

    if analysis.get("brain_insights"):
        console.print()
        console.print(Panel(
            analysis["brain_insights"],
            title="[bold]뇌 에이전트 판단[/bold]",
            border_style="magenta", padding=(1,2)
        ))

    if analysis.get("gpt_analysis"):
        console.print()
        console.print(Panel(
            analysis["gpt_analysis"],
            title="[bold]GPT-4o 시장 분석[/bold]",
            border_style="green", padding=(1,2)
        ))

    if analysis.get("claude_insights"):
        console.print()
        console.print(Panel(
            analysis["claude_insights"],
            title="[bold]전략 인사이트[/bold]",
            border_style="yellow", padding=(1,2)
        ))

    if analysis.get("sources"):
        console.print()
        tbl = Table(title="참고 출처", border_style="dim", show_header=True, header_style="bold")
        tbl.add_column("#", style="dim", width=3)
        tbl.add_column("제목", style="cyan", max_width=50)
        tbl.add_column("URL", style="dim", max_width=55)
        for i, s in enumerate(analysis["sources"], 1):
            tbl.add_row(str(i), s["title"][:50], s["url"][:55])
        console.print(tbl)

# ═══════════════════════════════════════════════════════
#  파일 저장
# ═══════════════════════════════════════════════════════
def save_results(topic: str, plan: dict, analysis: dict, results: list):
    console.print()
    console.print("[bold]저장 형식:[/bold]  [cyan]1[/cyan] Excel  [cyan]2[/cyan] PDF  "
                  "[cyan]3[/cyan] PPT  [cyan]4[/cyan] 전체  [cyan]0[/cyan] 건너뜀")
    choice = Prompt.ask("선택", choices=["0","1","2","3","4"], default="1")
    if choice == "0":
        return

    ts   = datetime.now().strftime("%Y%m%d_%H%M")
    safe = re.sub(r'[^\w가-힣]', '_', topic)[:30]
    base = f"리서치_{safe}_{ts}"
    saved = []

    if choice in ("1","4"):
        p = _save_excel(base, topic, analysis, results)
        if p: saved.append(p)

    if choice in ("2","4") and _PDF:
        p = _save_pdf(base, topic, analysis)
        if p: saved.append(p)
    elif choice == "2" and not _PDF:
        console.print("[yellow]reportlab 패키지 없음 — PDF 불가[/yellow]")

    if choice in ("3","4") and _PPT:
        p = _save_ppt(base, topic, analysis)
        if p: saved.append(p)
    elif choice == "3" and not _PPT:
        console.print("[yellow]python-pptx 패키지 없음 — PPT 불가[/yellow]")

    if saved:
        console.print()
        console.print(Panel(
            "\n".join(f"[green]OK[/green]  {s}" for s in saved),
            title="[bold]저장 완료[/bold]", border_style="green"
        ))

def _save_excel(base, topic, analysis, results):
    try:
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter

        per_url = analysis.get("per_url", [])
        needs   = analysis.get("needs", [])
        path    = os.path.join(SAVE_DIR, base+".xlsx")

        THIN = Side(style="thin", color="CCCCCC")
        def cell(ws, row, col, value, bold=False, bg=None, fc="000000",
                 wrap=True, size=10, align="left"):
            c = ws.cell(row=row, column=col, value=value)
            c.font      = Font(bold=bold, color=fc, size=size, name="맑은 고딕")
            c.alignment = Alignment(horizontal=align, vertical="top", wrap_text=wrap)
            c.border    = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
            if bg:
                c.fill = PatternFill("solid", fgColor=bg)
            return c

        wb = openpyxl.Workbook()

        ws = wb.active
        ws.title = "리서치 결과"

        fixed_cols  = [("No",5),("출처",8),("URL",32),("도메인",18),("페이지 제목",26),("한줄 요약",30)]
        all_headers = fixed_cols + [(n, 28) for n in needs] + [("헤딩 구조",28),("오류",14)]
        total_cols  = len(all_headers)

        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=total_cols)
        tc = ws.cell(row=1, column=1,
                     value=f"웹 리서치 — {topic}  |  {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        tc.font      = Font(bold=True, size=13, color="FFFFFF", name="맑은 고딕")
        tc.fill      = PatternFill("solid", fgColor="1F3864")
        tc.alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[1].height = 28

        for ci, (hdr, width) in enumerate(all_headers, 1):
            cell(ws, 2, ci, hdr, bold=True, bg="2F5496", fc="FFFFFF", align="center")
            ws.column_dimensions[get_column_letter(ci)].width = width
        ws.row_dimensions[2].height = 22

        for ri, data in enumerate(per_url, 1):
            er  = ri + 2
            bg  = "EBF0FA" if ri % 2 else "FFFFFF"
            gpt = data.get("gpt", {})
            ws.row_dimensions[er].height = 80
            src = "네이버" if "naver" in data.get("url","") else "웹"
            row_vals = [ri, src, data.get("url",""), data.get("domain",""),
                        data.get("page_title", data.get("title","")), gpt.get("한줄요약","")]
            for n in needs:
                row_vals.append(gpt.get(n, ""))
            row_vals += [data.get("headings",""), data.get("error","")]
            for ci, val in enumerate(row_vals, 1):
                cell(ws, er, ci, str(val) if val else "", bg=bg,
                     align="center" if ci == 1 else "left")

        ws.freeze_panes = "B3"

        ws2 = wb.create_sheet("종합 분석")
        ws2.column_dimensions["A"].width = 20
        ws2.column_dimensions["B"].width = 90

        def r2(row, label, text):
            ws2.cell(row=row, column=1, value=label).font = Font(bold=True, name="맑은 고딕", size=10)
            c = ws2.cell(row=row, column=2, value=text)
            c.alignment = Alignment(wrap_text=True, vertical="top")
            c.font = Font(name="맑은 고딕", size=10)
            ws2.row_dimensions[row].height = max(60, len(str(text))//3)

        r2(1, "리서치 주제",    topic)
        r2(2, "뇌 에이전트",   analysis.get("brain_insights","(미연동)"))
        r2(3, "GPT-4o 분석",   analysis.get("gpt_analysis",""))
        r2(4, "전략 인사이트", analysis.get("claude_insights",""))

        wb.save(path)
        return path
    except Exception as e:
        console.print(f"[red]Excel 오류: {e}[/red]")

def _save_pdf(base, topic, analysis):
    try:
        path = os.path.join(SAVE_DIR, base+".pdf")
        pdfmetrics.registerFont(UnicodeCIDFont('HYSMyeongJo-Medium'))
        ks = ParagraphStyle('k', fontName='HYSMyeongJo-Medium', fontSize=10, leading=16)
        kh = ParagraphStyle('h', fontName='HYSMyeongJo-Medium', fontSize=14, leading=20, spaceAfter=8)
        doc = SimpleDocTemplate(path, pagesize=A4)
        story = [Paragraph(f"리서치 리포트: {topic}", kh), Spacer(1,10)]
        if analysis.get("brain_insights"):
            story.append(Paragraph("뇌 에이전트 판단", kh))
            for ln in analysis["brain_insights"].split("\n"):
                if ln.strip(): story.append(Paragraph(ln, ks))
        story.append(Paragraph("GPT-4o 분석", kh))
        for ln in analysis.get("gpt_analysis","").split("\n"):
            if ln.strip(): story.append(Paragraph(ln, ks))
        story += [Spacer(1,10), Paragraph("전략 인사이트", kh)]
        for ln in analysis.get("claude_insights","").split("\n"):
            if ln.strip(): story.append(Paragraph(ln, ks))
        doc.build(story)
        return path
    except Exception as e:
        console.print(f"[red]PDF 오류: {e}[/red]")

def _save_ppt(base, topic, analysis):
    try:
        path = os.path.join(SAVE_DIR, base+".pptx")
        prs = Presentation()
        def add_slide(title_text, body_text):
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = title_text
            sl.placeholders[1].text = body_text[:900]
        sl0 = prs.slides.add_slide(prs.slide_layouts[0])
        sl0.shapes.title.text = f"리서치 리포트: {topic}"
        sl0.placeholders[1].text = datetime.now().strftime("%Y-%m-%d")
        if analysis.get("brain_insights"):
            add_slide("뇌 에이전트 판단", analysis["brain_insights"][:900])
        add_slide("GPT-4o 분석",   analysis.get("gpt_analysis","")[:900])
        add_slide("전략 인사이트", analysis.get("claude_insights","")[:900])
        prs.save(path)
        return path
    except Exception as e:
        console.print(f"[red]PPT 오류: {e}[/red]")

# ═══════════════════════════════════════════════════════
#  오케스트레이터 피드백 수집
# ═══════════════════════════════════════════════════════
def _collect_feedback(topic: str, plan: dict, analysis: dict):
    """리서치 결과에 대한 사용자 평점 수집 → 오케스트레이터 학습 DB 저장"""
    console.print()
    console.print(Rule("[dim]오케스트레이터 학습[/dim]", style="dim"))

    # 사용 모델 표시
    if _session_models:
        model_summary = "  ".join(f"[dim]{cat}:[/dim] [cyan]{m}[/cyan]"
                                   for cat, m in _session_models.items())
        console.print(f"[dim]이번 세션 모델: {model_summary}[/dim]")

    score_raw = Prompt.ask(
        "[dim]이번 리서치 만족도 (1~5, 건너뛰려면 엔터)[/dim]",
        default=""
    ).strip()
    if not score_raw or not score_raw.isdigit():
        return

    score = max(1, min(5, int(score_raw)))

    needs = analysis.get("needs", [])
    useful_needs, bad_needs = [], []

    if needs:
        console.print(f"[dim]추출 항목: {', '.join(f'{i+1}.{n}' for i, n in enumerate(needs))}[/dim]")
        useful_raw = Prompt.ask("[dim]유용했던 항목 번호 (쉼표, 없으면 엔터)[/dim]", default="").strip()
        bad_raw    = Prompt.ask("[dim]불필요했던 항목 번호 (쉼표, 없으면 엔터)[/dim]", default="").strip()

        def _parse_nums(raw, items):
            if not raw:
                return []
            return [items[int(x.strip()) - 1]
                    for x in raw.split(",")
                    if x.strip().isdigit() and 0 < int(x.strip()) <= len(items)]

        useful_needs = _parse_nums(useful_raw, needs)
        bad_needs    = _parse_nums(bad_raw, needs)

    comment = Prompt.ask("[dim]한마디 코멘트 (없으면 엔터)[/dim]", default="").strip()

    orch.record_feedback(
        topic=topic, plan=plan, analysis=analysis,
        score=score, useful_needs=useful_needs, bad_needs=bad_needs,
        models_used=dict(_session_models), comment=comment
    )
    console.print(f"[dim]피드백 저장 완료 (점수: {score}/5)[/dim]")


def _push_feedback(topic: str, plan: dict, analysis: dict, conversation: list):
    """
    세션 대화 + 피드백 → 암호화 로컬 저장 + GitHub Issue 자동 전송.
    오케스트레이터 피드백(_collect_feedback)에서 이미 score/needs 수집했으면 재사용.
    """
    # _collect_feedback에서 저장한 마지막 피드백 DB에서 꺼내기
    score, useful_needs, bad_needs = 0, [], []
    try:
        import sqlite3
        con = sqlite3.connect(orch.DB_PATH)
        row = con.execute(
            "SELECT score, useful_needs, bad_needs FROM feedback_records ORDER BY id DESC LIMIT 1"
        ).fetchone()
        con.close()
        if row:
            score        = row[0] or 0
            useful_needs = json.loads(row[1] or "[]")
            bad_needs    = json.loads(row[2] or "[]")
    except Exception:
        pass

    with console.status("[dim]인사이트 추출 중...[/dim]"):
        result = fc.process_session(
            topic=topic,
            research_type=plan.get("research_type", "general"),
            plan=plan, analysis=analysis,
            conversation=conversation,
            score=score,
            useful_needs=useful_needs, bad_needs=bad_needs,
            models_used=dict(_session_models),
        )

    if result.get("log_path"):
        console.print(f"[dim]대화 암호화 저장: {os.path.basename(result['log_path'])}[/dim]")

    if result.get("issue_url"):
        console.print(f"[dim]GitHub 이슈 생성: {result['issue_url']}[/dim]")
    elif result.get("insights"):
        # GitHub 실패해도 로컬에는 저장됨
        console.print("[dim]인사이트 추출 완료 (GitHub 전송 실패 — 로컬 저장됨)[/dim]")


# ═══════════════════════════════════════════════════════
#  인테이크 채팅 — GPT와 대화로 조사 주제 구체화
# ═══════════════════════════════════════════════════════
def intake_chat(first_input: str):
    """GPT와 자연스러운 대화로 조사 주제와 깊이를 확정 후 (topic, depth) 반환"""
    if not oai:
        # API 없으면 입력값 그대로 사용
        return first_input, "일반 조사"

    system = """당신은 웹 리서치 어시스턴트의 접수 담당입니다.
사용자가 조사를 요청하면 자연스럽게 대화해서 주제를 구체화하세요.

중요 규칙:
1. 반드시 JSON 형식으로만 응답하세요. 일반 텍스트 금지.
2. 질문은 한 번에 하나만.
3. 주제가 충분히 명확하면 (1~2턴이면 충분) 바로 확정.
4. 조사와 무관한 잡담이면 친절하게 조사 주제로 유도.

응답 형식:

대화 중일 때:
{"ready": false, "message": "여기에 실제 질문이나 안내 문구를 작성"}

예시: {"ready": false, "message": "어떤 업종의 경쟁사를 조사할까요?"}
예시: {"ready": false, "message": "안녕하세요! 어떤 주제를 조사해드릴까요?"}

주제 확정됐을 때:
{"ready": true, "topic": "확정된 조사 주제", "depth": "빠른 조사 또는 일반 조사 또는 심층 조사", "message": "네, [주제] 조사를 시작하겠습니다!"}

종료 요청 시:
{"ready": false, "quit": true, "message": ""}

깊이 기준: 빠른=가볍게 궁금한 것, 일반=업무용, 심층=중요한 의사결정"""

    history = [{"role": "user", "content": first_input}]

    for _ in range(6):  # 최대 6턴
        try:
            resp = oai.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": system}] + history,
                max_tokens=300, temperature=0.5
            )
            raw = resp.choices[0].message.content.strip()
        except Exception:
            return first_input, "일반 조사"

        # JSON 파싱 시도
        parsed = None
        try:
            m = re.search(r'\{[\s\S]+\}', raw)
            if m:
                parsed = json.loads(m.group())
        except Exception:
            pass

        if not parsed:
            # JSON 없으면 그냥 텍스트 출력 후 사용자 입력
            console.print(f"\n[bold cyan]어시스턴트[/bold cyan]: {raw}")
            user_in = Prompt.ask("[bold cyan]→[/bold cyan]").strip()
            if not user_in or user_in.lower() == "q":
                return None
            history.append({"role": "assistant", "content": raw})
            history.append({"role": "user", "content": user_in})
            continue

        if parsed.get("quit"):
            return None

        msg = parsed.get("message", "")

        if parsed.get("ready"):
            topic = parsed.get("topic", first_input)
            depth = parsed.get("depth", "일반 조사")
            if msg:
                console.print(f"\n[bold cyan]어시스턴트[/bold cyan]: {msg}")
            console.print()
            return topic, depth
        else:
            # 아직 대화 중
            if msg:
                console.print(f"\n[bold cyan]어시스턴트[/bold cyan]: {msg}")
            user_in = Prompt.ask("[bold cyan]→[/bold cyan]").strip()
            if not user_in or user_in.lower() == "q":
                return None
            history.append({"role": "assistant", "content": raw})
            history.append({"role": "user", "content": user_in})

    # 6턴 초과 시 첫 입력 그대로 사용
    return first_input, "일반 조사"


# ═══════════════════════════════════════════════════════
#  채팅 모드 — 리서치 결과 기반 대화
# ═══════════════════════════════════════════════════════
def chat_mode(topic: str, plan: dict, analysis: dict, results: list) -> list:
    """리서치 결과를 컨텍스트로 Claude와 대화. 대화 히스토리 반환."""
    console.print()
    console.print(Rule("[bold magenta]결과 기반 대화 모드[/bold magenta]", style="magenta"))
    console.print("[dim]수집된 데이터를 바탕으로 추가 질문하세요. 종료: q 또는 엔터[/dim]")

    # 컨텍스트 구성
    per_url_summary = "\n".join(
        f"- [{d.get('domain','')}] {d.get('gpt',{}).get('한줄요약','')}"
        for d in analysis.get("per_url", [])[:10]
    )
    ctx = (
        f"리서치 주제: {topic}\n"
        f"조사 유형: {plan.get('research_type','')}\n\n"
        f"수집된 사이트별 요약:\n{per_url_summary}\n\n"
        f"GPT-4o 시장 분석:\n{analysis.get('gpt_analysis','')[:1500]}\n\n"
        f"전략 인사이트:\n{analysis.get('claude_insights','')[:800]}\n\n"
        f"뇌 에이전트 판단:\n{analysis.get('brain_insights','(미연동)')}"
    )
    history = []  # 대화 히스토리

    while True:
        console.print()
        q = Prompt.ask("[bold magenta]질문[/bold magenta]").strip()
        if not q or q.lower() == "q":
            break

        history.append({"role": "user", "content": q})
        msgs = [
            {"role": "system", "content":
                f"당신은 아래 웹 리서치 결과를 완전히 숙지한 전문 마케팅 분석가입니다.\n"
                f"실제 수집된 데이터 기반으로만 답하고, 모르면 솔직히 말하세요.\n\n"
                f"=== 리서치 컨텍스트 ===\n{ctx}"},
        ] + history

        with console.status("[dim]분석 중...[/dim]"):
            try:
                if claude_c:
                    r = claude_c.messages.create(
                        model="claude-opus-4-6", max_tokens=1500,
                        system=msgs[0]["content"],
                        messages=history
                    )
                    answer = r.content[0].text.strip()
                elif oai:
                    r = oai.chat.completions.create(
                        model="gpt-4o", messages=msgs, max_tokens=1500
                    )
                    answer = r.choices[0].message.content.strip()
                else:
                    answer = "API 키가 없습니다."
            except Exception as e:
                answer = f"오류: {e}"

        history.append({"role": "assistant", "content": answer})
        console.print()
        console.print(Panel(answer, border_style="magenta", padding=(0,2)))

    console.print("[dim]대화 모드 종료[/dim]")
    return history


# ═══════════════════════════════════════════════════════
#  메인
# ═══════════════════════════════════════════════════════
def main():
    console.clear()
    console.print(Panel(
        "[bold cyan]웹 리서치 어시스턴트[/bold cyan]\n"
        "[dim]DeepSeek R1 · GPT-4o · Claude · 뇌 에이전트 기반 자동 웹조사[/dim]",
        border_style="bright_blue", padding=(1,6)
    ))
    console.print(f"\n[dim]결과 저장 폴더: {SAVE_DIR}[/dim]")

    # 오케스트레이터 백그라운드 재평가 (7일 이상 된 카테고리)
    _start_background_evaluation()

    # 뇌 에이전트 상태 출력 — HTTP API 먼저, 로컬 폴백
    _brain_ok = False
    try:
        _ping = requests.get(
            _BRAIN_AGENT_URL.replace("/api/research", "/"),
            timeout=5
        )
        _brain_ok = _ping.status_code < 500
    except Exception:
        _brain_ok = bool(_get_brain_agent())
    if _brain_ok:
        console.print("[dim]뇌 에이전트: 연동됨[/dim]")
    else:
        console.print("[dim]뇌 에이전트: 미연동 (Claude로 대체)[/dim]")

    while True:
        console.print()
        console.print(Rule(style="dim"))
        console.print("[bold]무엇이 궁금하세요?[/bold]")
        console.print("[dim]자유롭게 말씀해 주세요. 대화하면서 조사 방향을 잡아드립니다.[/dim]")
        console.print("[dim]종료: q[/dim]")

        first_input = Prompt.ask("\n[bold cyan]→[/bold cyan]").strip()
        if not first_input or first_input.lower() == "q":
            console.print("\n[dim]프로그램을 종료합니다.[/dim]")
            break

        # 세션 모델 트래커 초기화
        _session_models.clear()

        # GPT 대화로 주제 구체화
        result = intake_chat(first_input)
        if not result:
            continue
        topic, depth = result

        console.print()
        console.print(Rule("[cyan]리서치 시작[/cyan]", style="bright_blue"))

        # 1. 플랜
        with console.status("[cyan]DeepSeek R1 — 조사 유형 판단 + 계획 수립 중...[/cyan]"):
            plan = make_plan(topic, depth)
            plan["_depth"] = depth
        console.print(f"[green]OK[/green] 검색 쿼리 {len(plan['queries'])}개 생성")

        # needs 사용자 확인
        needs = plan.get("needs", [])
        console.print()
        console.print("[bold]추출 항목 확인[/bold]  (DeepSeek 추천):")
        for i, n in enumerate(needs, 1):
            console.print(f"  [cyan]{i}[/cyan] {n}")
        console.print("[dim]Enter=전체 사용  /  번호(예: 1,3,5)  /  직접입력(쉼표 구분)[/dim]")
        needs_raw = Prompt.ask("선택", default="").strip()
        if not needs_raw:
            pass
        elif re.match(r'^[\d,\s]+$', needs_raw):
            idx = [int(x.strip()) - 1 for x in needs_raw.split(",") if x.strip().isdigit()]
            selected = [needs[i] for i in idx if 0 <= i < len(needs)]
            if selected:
                plan["needs"] = selected
        else:
            plan["needs"] = [n.strip() for n in needs_raw.replace("，", ",").split(",") if n.strip()]

        # 2. 수집
        console.print("[cyan]웹 검색 및 페이지 수집 중...[/cyan]")
        results = run_research(plan)
        console.print(f"[green]OK[/green] {len(results)}개 페이지 수집 완료")

        # 3. 분석
        console.print("[cyan]GPT 구조화 추출 + 종합 분석 중...[/cyan]")
        analysis = analyze(topic, plan, results)
        console.print("[green]OK[/green] 분석 완료")

        # 4. 출력
        display_results(topic, plan, analysis)

        # 5. 저장
        save_results(topic, plan, analysis, results)

        # 5.5 오케스트레이터 피드백 수집
        if _ORCH_ENABLED:
            _collect_feedback(topic, plan, analysis)

        # 6. 채팅 모드 (대화 히스토리 수집)
        conversation = chat_mode(topic, plan, analysis, results)

        # 7. 피드백 → 암호화 저장 + GitHub 이슈
        if _FC_ENABLED:
            _push_feedback(topic, plan, analysis, conversation)

        # 7. 계속?
        console.print()
        if not Confirm.ask("[bold]새 주제를 조사할까요?[/bold]", default=True):
            console.print("\n[dim]프로그램을 종료합니다.[/dim]")
            break

if __name__ == "__main__":
    # MAESTRO로 통합 — app_local.py는 리서치 파이프라인 모듈로 동작
    try:
        import maestro
        maestro.main()
    except ImportError:
        # maestro.py 없으면 기존 방식으로 실행 (하위 호환)
        main()
