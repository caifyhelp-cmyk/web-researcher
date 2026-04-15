# -*- coding: utf-8 -*-
"""웹 리서치 어시스턴트 v3 — 진짜 맞춤형 LLM (피드백 = 시스템 프롬프트 + 파이프라인 규칙)"""

# PDF / PPT 내보내기
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    _PDF_OK = True
except ImportError:
    _PDF_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    _PPT_OK = True
except ImportError:
    _PPT_OK = False

import streamlit as st
import pandas as pd
import json, io, os, sys, re, time, random, hashlib
from datetime import datetime
import requests as http_req
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from openai import OpenAI
from anthropic import Anthropic
from web_researcher import (
    OPENAI_API_KEY,
    gpt_suggest_needs, gpt_suggest_context,
    search_all_engines, discover_internal_pages,
    scrape_page, analyze_with_gpt,
    make_driver, is_login_wall, validate_domain,
)

# ─────────────────────────────────────────────
#  상수
# ─────────────────────────────────────────────
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")
DEEPSEEK_API_KEY  = os.getenv("DEEPSEEK_API_KEY",  "")
GROK_API_KEY      = os.getenv("GROK_API_KEY",      "")
BRAIN_AGENT_URL   = "https://brain-agent-v9wl.onrender.com"
BASE_DIR   = os.path.dirname(os.path.abspath(__file__))
DATA_DIR   = os.path.join(BASE_DIR, "data")
USERS_FILE = os.path.join(DATA_DIR, "users.json")
os.makedirs(DATA_DIR, exist_ok=True)

# ─────────────────────────────────────────────
#  GitHub 영구 저장소 (재시작해도 데이터 유지)
# ─────────────────────────────────────────────
import base64

_GH_TOKEN = os.getenv("GITHUB_DATA_TOKEN", "")
_GH_REPO  = "caifyhelp-cmyk/web-researcher-data"
_GH_API   = "https://api.github.com"
_GH_HEADS = lambda: {
    "Authorization": f"token {_GH_TOKEN}",
    "Accept": "application/vnd.github+json",
    "Content-Type": "application/json",
}

def _gh_read(path: str):
    """GitHub에서 JSON 파일 읽기. (dict, sha) 반환"""
    if not _GH_TOKEN:
        return {}, None
    try:
        r = http_req.get(f"{_GH_API}/repos/{_GH_REPO}/contents/{path}",
                         headers=_GH_HEADS(), timeout=10)
        if r.status_code == 404:
            return {}, None
        data = r.json()
        content = json.loads(base64.b64decode(data["content"]).decode("utf-8"))
        return content, data["sha"]
    except Exception:
        return {}, None

def _gh_write(path: str, content: dict, sha=None, msg="update"):
    """GitHub에 JSON 파일 쓰기"""
    if not _GH_TOKEN:
        return False
    try:
        encoded = base64.b64encode(json.dumps(content, ensure_ascii=False, indent=2).encode("utf-8")).decode()
        body = {"message": msg, "content": encoded}
        if sha:
            body["sha"] = sha
        r = http_req.put(f"{_GH_API}/repos/{_GH_REPO}/contents/{path}",
                         headers=_GH_HEADS(), json=body, timeout=15)
        return r.status_code in (200, 201)
    except Exception:
        return False



MARKETING_KW = [
    "랜딩","광고","마케팅","전환","경쟁사","cpc","seo","콘텐츠","sns",
    "유튜브","쇼츠","브랜딩","캠페인","리드","퍼널","전략","홍보",
    "프로모션","인플루언서","키워드","클릭","노출","분석","페이지",
]

# ─────────────────────────────────────────────
#  페이지 설정 + CSS
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="웹 리서치 어시스턴트",
    page_icon="🔍", layout="wide",
    initial_sidebar_state="expanded",
)
st.markdown("""
<style>
  .stApp { background:#0f1117; }
  .block-container { max-width:880px; padding-top:1.5rem; }
  #MainMenu,footer,header { visibility:hidden; }

  .app-hdr {
    background:linear-gradient(135deg,#1a1f36,#16213e);
    border:1px solid #2a2f4e; border-radius:14px;
    padding:22px 32px; margin-bottom:22px;
  }
  .app-hdr h1 { font-size:22px; font-weight:700; color:#e2e8f0; margin:0; }
  .app-hdr p  { font-size:13px; color:#64748b; margin:3px 0 0; }

  .step-lbl {
    font-size:10px; font-weight:700; letter-spacing:1.5px;
    color:#6366f1; text-transform:uppercase; margin-bottom:4px;
  }
  .profile-box {
    background:#1a1f36; border:1px solid #7c3aed;
    border-radius:10px; padding:12px 16px; margin-top:8px;
    font-size:12px; color:#c4b5fd; line-height:1.7;
  }
  .rule-chip {
    display:inline-block; background:#1e1b4b; color:#a78bfa;
    border:1px solid #4c1d95; border-radius:12px;
    padding:2px 10px; font-size:11px; margin:2px;
  }
  .ins-title { font-size:13px; font-weight:700; color:#818cf8; margin-bottom:8px; }
  .brain-badge {
    display:inline-block; background:#2d1b69; color:#a78bfa;
    border:1px solid #7c3aed; border-radius:6px;
    padding:2px 10px; font-size:11px; font-weight:600; margin-bottom:10px;
  }
  .res-card {
    background:#1a1f36; border:1px solid #2a2f4e;
    border-radius:10px; padding:14px 18px; margin-bottom:8px;
  }
  .res-domain { font-size:13px; font-weight:600; color:#818cf8; }
  .res-url    { font-size:11px; color:#475569; }
  .res-sum    { font-size:13px; color:#94a3b8; margin-top:6px; line-height:1.6; }
  .change-item {
    background:#052e16; border:1px solid #166534; border-radius:8px;
    padding:8px 14px; margin-bottom:6px; font-size:13px; color:#86efac;
  }
  .stButton>button {
    background:#4338ca!important; color:#fff!important;
    border:none!important; border-radius:8px!important;
    font-weight:600!important; padding:10px 0!important; width:100%;
  }
  .stButton>button:hover { background:#3730a3!important; }
  div[data-testid="stSidebarContent"] { background:#0d1117; }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────
#  클라이언트
# ─────────────────────────────────────────────
oai_client = OpenAI(api_key=OPENAI_API_KEY)
claude      = Anthropic(api_key=ANTHROPIC_API_KEY)
deepseek    = OpenAI(api_key=DEEPSEEK_API_KEY, base_url="https://api.deepseek.com")
grok        = OpenAI(api_key=GROK_API_KEY,     base_url="https://api.x.ai/v1")

# ─────────────────────────────────────────────
#  인증
# ─────────────────────────────────────────────
def _hash_pw(pw, salt=None):
    if salt is None: salt = os.urandom(16)
    key = hashlib.pbkdf2_hmac("sha256", pw.encode(), salt, 200_000)
    return salt.hex() + ":" + key.hex()

def _verify_pw(pw, stored):
    try:
        salt, kh = stored.split(":")
        k = hashlib.pbkdf2_hmac("sha256", pw.encode(), bytes.fromhex(salt), 200_000)
        return k.hex() == kh
    except Exception:
        return False

@st.cache_data(ttl=60, show_spinner=False)
def load_users():
    if _GH_TOKEN:
        data, _ = _gh_read("users.json")
        return data
    return json.load(open(USERS_FILE, encoding="utf-8")) if os.path.exists(USERS_FILE) else {}

def save_users(u):
    load_users.clear()
    if _GH_TOKEN:
        _, sha = _gh_read("users.json")
        _gh_write("users.json", u, sha, "save users")
        return
    with open(USERS_FILE, "w", encoding="utf-8") as f:
        json.dump(u, f, ensure_ascii=False, indent=2)

def authenticate(username, pw):
    u = load_users().get(username)
    return u if u and _verify_pw(pw, u["password_hash"]) else None

def register_user(name, pw):
    username = name.strip()
    if len(username) < 1: return False, "이름을 입력하세요"
    if len(pw) < 4:       return False, "암호는 4자 이상"
    users = load_users()
    if username in users: return False, "이미 사용 중인 아이디"
    users[username] = {"name": name, "password_hash": _hash_pw(pw),
                       "created_at": datetime.now().isoformat()}
    save_users(users)
    return True, ""

# ─────────────────────────────────────────────
#  유저 프로필 (맞춤 LLM의 핵심)
# ─────────────────────────────────────────────
def user_dir(username):
    d = os.path.join(DATA_DIR, username)
    os.makedirs(d, exist_ok=True)
    return d

@st.cache_data(ttl=300, show_spinner=False)
def load_profile(username) -> dict:
    if _GH_TOKEN:
        data, _ = _gh_read(f"profiles/{username}.json")
        return data if data else {"system_prompt": "", "rules": {}, "history": []}
    d = os.path.join(DATA_DIR, username)
    p = os.path.join(d, "profile.json")
    if os.path.exists(p):
        try:
            return json.load(open(p, encoding="utf-8"))
        except Exception:
            pass
    return {"system_prompt": "", "rules": {}, "history": []}
def save_profile(username, profile):
    load_profile.clear()
    if _GH_TOKEN:
        _, sha = _gh_read(f"profiles/{username}.json")
        _gh_write(f"profiles/{username}.json", profile, sha, f"save profile {username}")
        return
    d = os.path.join(DATA_DIR, username)
    os.makedirs(d, exist_ok=True)
    with open(os.path.join(d, "profile.json"), "w", encoding="utf-8") as f:
        json.dump(profile, f, ensure_ascii=False, indent=2)


@st.cache_data(ttl=300, show_spinner=False)
def load_history(username) -> list:
    if _GH_TOKEN:
        hist, _ = _gh_read(f"history/{username}.json")
        return hist if isinstance(hist, list) else []
    d = os.path.join(DATA_DIR, username)
    hf = os.path.join(d, "history.json")
    if os.path.exists(hf):
        try:
            return json.load(open(hf, encoding="utf-8"))
        except Exception:
            pass
    return []


def save_history(username, entry):
    load_history.clear()
    if _GH_TOKEN:
        hist, sha = _gh_read(f"history/{username}.json")
        if not isinstance(hist, list):
            hist = []
        hist.append(entry)
        hist = hist[-50:]
        _gh_write(f"history/{username}.json", hist, sha, f"history {username}")
        return
    d = os.path.join(DATA_DIR, username)
    os.makedirs(d, exist_ok=True)
    hf = os.path.join(d, "history.json")
    hist = json.load(open(hf, encoding="utf-8")) if os.path.exists(hf) else []
    hist.append(entry)
    with open(hf, "w", encoding="utf-8") as f:
        json.dump(hist[-50:], f, ensure_ascii=False, indent=2)


def apply_user_rules(config: dict, plan: dict, rules: dict) -> tuple:
    """사용자 규칙을 config/plan에 직접 적용 — 파이프라인 실질 변경"""
    config = dict(config)
    plan   = dict(plan)

    # 수집 개수
    if rules.get("min_count") and config.get("count", 0) < rules["min_count"]:
        config["count"] = rules["min_count"]
    if rules.get("count_multiplier", 1.0) > 1.0:
        config["count"] = min(100, int(config["count"] * rules["count_multiplier"]))

    # 필터 강도
    if rules.get("force_official_only") or rules.get("exclude_blogs"):
        plan["strict_filter"] = "high"
    elif rules.get("force_strict_filter") and rules["force_strict_filter"] != "null":
        plan["strict_filter"] = rules["force_strict_filter"]

    # 도메인 제외
    extra = list(rules.get("extra_excluded_domains", []))
    if extra:
        plan["exclude_extra"] = list(plan.get("exclude_extra", [])) + extra

    return config, plan


def get_user_system(username: str) -> str:
    """Claude 모든 호출에 주입되는 개인 시스템 프롬프트"""
    sp = load_profile(username).get("system_prompt", "")
    return f"[이 사용자 특성 — 반드시 반영]\n{sp}" if sp else ""


# ─────────────────────────────────────────────
#  Claude LLM 함수 (시스템 프롬프트 주입)
# ─────────────────────────────────────────────
def build_plan_claude(config: dict, username: str) -> dict:
    """
    리서치 플랜 수립
    1차: DeepSeek deepseek-chat (추론 특화, 저렴)
    fallback: Claude Sonnet 4.6
    개인 시스템 프롬프트 + 규칙 힌트 반영
    """
    user_system = get_user_system(username)
    rules = load_profile(username).get("rules", {})

    hints = []
    if rules.get("force_official_only"): hints.append("공식 기관 사이트만 수집")
    if rules.get("exclude_blogs"):       hints.append("블로그·개인 사이트 제외")
    if rules.get("min_count"):           hints.append(f"최소 {rules['min_count']}개 수집")
    if rules.get("query_style"):         hints.append(f"쿼리 스타일: {rules['query_style']}")
    hints_block = "\n".join(f"- {h}" for h in hints) if hints else ""

    sys_parts = []
    if user_system: sys_parts.append(user_system)
    sys_parts.append("당신은 한국 웹 리서치 플랜 전문가입니다. 유효한 JSON만 출력하세요.")
    system_msg = "\n\n".join(sys_parts)

    kw = config["keyword"]
    prompt = "\n".join(filter(None, [
        f'조사: "{kw}"',
        f'항목: {', '.join(config.get("needs", []))}',
        f'맥락: {config.get("context","없음")}',
        "",
        ("━━ 이 사용자 맞춤 요구사항 (반드시 반영) ━━\n" + hints_block) if hints_block else "",
        "",
        "━━ 조사 유형 판단 ━━",
        "A) 특정 기관·사이트 집중 분석",
        "B) 시장·경쟁사 비교",
        "C) 특정 정보·데이터 수집",
        "",
        "━━ 쿼리 생성 규칙 ━━",
        "- 실제 고객이 네이버/구글에 입력할 자연스러운 한국어 검색어",
        "- '분析', '전환율', '최적화' 같은 마케터 내부 용어 절대 금지",
        "- 목적에 맞는 형태: [키워드] 기관, [키워드] 협회, [키워드] 추천",
        "",
        "━━ 반드시 아래 JSON 형식으로만 출력 ━━",
        '{"type":"A or B or C","reason":"한 줄","target_domain":"","naver_queries":["q1","q2","q3","q4"],"google_queries":["q1","q2","q3","q4"],"entity_queries":[],"exclude_extra":[],"strict_filter":"high or medium or low","is_bulk":false}',
    ]))

    # 1차: DeepSeek (추론 특화, 저렴)
    try:
        resp = deepseek.chat.completions.create(
            model="deepseek-reasoner",
            max_tokens=800,
            messages=[
                {"role": "system", "content": system_msg},
                {"role": "user",   "content": prompt},
            ],
        )
        text = resp.choices[0].message.content
        m = re.search(r'\{[\s\S]+\}', text)
        if m:
            result = json.loads(m.group())
            result["_llm"] = "deepseek-chat"
            return result
    except Exception as e:
        print(f"[플랜/DeepSeek] {e}")

    # fallback: Claude Sonnet
    try:
        resp = claude.messages.create(
            model="claude-sonnet-4-6", max_tokens=800,
            system=user_system or "당신은 한국 웹 리서치 플랜 전문가입니다.",
            messages=[{"role": "user", "content": prompt}],
        )
        m = re.search(r'\{[\s\S]+\}', resp.content[0].text)
        if m:
            result = json.loads(m.group())
            result["_llm"] = "claude-sonnet-4-6"
            return result
    except Exception as e:
        print(f"[플랜/Claude] {e}")

    return {
        "type": "B", "reason": "기본값", "target_domain": "",
        "naver_queries":  [kw, f"{kw} 기관", f"{kw} 추천", f"{kw} 비교"],
        "google_queries": [kw, f"best {kw}", f"{kw} comparison", f"{kw} review"],
        "entity_queries": [], "exclude_extra": [], "strict_filter": "medium",
        "is_bulk": False, "_llm": "fallback",
    }

def generate_insights(all_results: list, config: dict, username: str) -> str:
    """Claude Opus 4.6 — 개인 시스템 프롬프트 반영한 종합 인사이트"""
    user_system = get_user_system(username)
    lines = "\n".join(
        f"- {r.get('domain','')}: {r.get('gpt',{}).get('한줄요약','')}"
        for r in all_results
    )
    try:
        resp = claude.messages.create(
            model="claude-opus-4-6",
            max_tokens=2000,
            system=user_system or "당신은 실무 중심의 웹 리서치 전문가입니다.",
            messages=[{"role": "user", "content": f"""조사 주제: "{config['keyword']}"
수집 항목: {', '.join(config.get('needs',[]))}
수집 {len(all_results)}개:
{lines}

마크다운으로 (한국어, 실무 중심):
## 핵심 인사이트
## 공통점 · 패턴
## 차별화 포인트
## 실무 활용 제안"""}],
        )
        return resp.content[0].text.strip()
    except Exception as e:
        print(f"[인사이트] {e}")
        return "인사이트 생성 실패"


def get_grok_realtime(config: dict, insights: str) -> str:
    """
    Grok (xAI grok-3) — 실시간 시장 신호 보완
    Claude Opus 인사이트 이후 최신 동향/뉴스 관점 추가
    """
    kw  = config.get("keyword", "")
    ctx = config.get("context", "")
    try:
        resp = grok.chat.completions.create(
            model="grok-3-mini-fast",
            max_tokens=600,
            messages=[
                {
                    "role": "system",
                    "content": "당신은 최신 시장 동향과 실시간 뉴스에 밝은 한국 비즈니스 애널리스트입니다.",
                },
                {
                    "role": "user",
                    "content": (
                        f'조사 주제: "{kw}" (맥락: {ctx})\n\n'
                        f"기존 인사이트 요약:\n{insights[:600]}\n\n"
                        "최근 6개월 내 실시간 동향과 뉴스 관점에서 보완해주세요.\n\n"
                        "## 실시간 시장 신호\n"
                        "- 최근 주목할 만한 동향 3~5가지 (bullet)\n\n"
                        "## 놓칠 수 있는 변수\n"
                        "- 기존 분석에서 빠졌을 수 있는 최신 리스크/기회 2~3가지"
                    ),
                },
            ],
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        print(f"[Grok] {e}")
        return ""

# ─────────────────────────────────────────────
#  뇌 에이전트
# ─────────────────────────────────────────────
def is_marketing(config):
    text = " ".join([config.get("keyword",""), config.get("context",""),
                     *config.get("needs",[])]).lower()
    return any(k in text for k in MARKETING_KW)

def call_brain_agent(config, all_results, insights):
    findings = "\n".join(
        f"- {r.get('domain','')}: {r.get('gpt',{}).get('한줄요약','')}"
        for r in all_results[:8]
    )
    try:
        resp = http_req.post(f"{BRAIN_AGENT_URL}/api/ask", json={
            "industry":     config.get("keyword",""),
            "company_size": "중소",
            "challenge":    f"{config.get('keyword','')} 마케팅 전략 수립",
            "context":      f"리서치:\n{findings}\n\n인사이트:\n{insights[:500]}",
            "goal":         config.get("context","시장 이해 및 전략 도출"),
            "assets":       ", ".join(config.get("needs",[])),
            "constraints":  "",
        }, timeout=40)
        if resp.status_code == 200:
            d = resp.json()
            if d.get("status") == "ok":
                return d.get("result",{}), d.get("fired_patterns",[])
    except Exception as e:
        print(f"[뇌 에이전트] {e}")
    return None, []


# ─────────────────────────────────────────────
#  데이터 / xlsx
# ─────────────────────────────────────────────
def results_to_df(all_results, config):
    rows = []
    for i, data in enumerate(all_results, 1):
        g=data.get("gpt",{}); p=data.get("page",{})
        row={"No":i,"URL":data.get("url",""),"도메인":data.get("domain",""),
             "출처":data.get("source",""),"페이지 제목":p.get("page_title",""),
             "한줄 요약":g.get("한줄요약","")}
        for n in config.get("needs",[]): row[n]=g.get(n,"")
        rows.append(row)
    return pd.DataFrame(rows)


def make_pdf_bytes(all_results, config, research_type, insights, brain_result, grok_signal=""):
    """reportlab CJK 폰트로 한글 PDF 생성"""
    if not _PDF_OK:
        return None
    buf = io.BytesIO()
    try:
        pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))
        KR_FONT = "HYSMyeongJo-Medium"
    except Exception:
        KR_FONT = "Helvetica"
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=20*mm, bottomMargin=20*mm)
    story = []
    styles = getSampleStyleSheet()
    def ps(name, **kw):
        return ParagraphStyle(name, parent=styles["Normal"], fontName=KR_FONT, **kw)
    title_st   = ps("T", fontSize=18, spaceAfter=6, textColor=colors.HexColor("#6366f1"), leading=24)
    h1_st      = ps("H1", fontSize=13, spaceAfter=4, textColor=colors.HexColor("#334155"), leading=18, spaceBefore=10)
    body_st    = ps("B", fontSize=9, spaceAfter=4, leading=14, textColor=colors.HexColor("#1e293b"))
    caption_st = ps("C", fontSize=8, spaceAfter=2, textColor=colors.HexColor("#64748b"), leading=12)
    def safe(text):
        if not text: return ""
        return str(text).replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    kw = config.get("keyword", "")
    story.append(Paragraph(f"리서치 보고서: {safe(kw)}", title_st))
    story.append(Paragraph(f"타입: {research_type}  |  수집: {len(all_results)}개  |  생성: {datetime.now().strftime('%Y-%m-%d %H:%M')}", caption_st))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#6366f1"), spaceAfter=8))
    if insights:
        story.append(Paragraph("종합 인사이트 (Claude Opus 4.6)", h1_st))
        for line in insights.split(chr(10)):
            if line.strip():
                story.append(Paragraph(safe(line.strip()), body_st))
        story.append(Spacer(1, 6*mm))
    if grok_signal:
        story.append(Paragraph("실시간 시장 신호 (Grok)", h1_st))
        for line in grok_signal.split(chr(10)):
            if line.strip():
                story.append(Paragraph(safe(line.strip()), body_st))
        story.append(Spacer(1, 6*mm))
    if brain_result:
        story.append(Paragraph("뇌 에이전트 마케팅 전략", h1_st))
        j = brain_result.get("judgment","")
        a = brain_result.get("action","")
        if j: story.append(Paragraph(f"판단: {safe(j)}", body_st))
        if a: story.append(Paragraph(f"액션: {safe(a)}", body_st))
        story.append(Spacer(1, 6*mm))
    if all_results:
        story.append(Paragraph("수집 사이트 목록", h1_st))
        for i, r in enumerate(all_results, 1):
            g = r.get("gpt", {})
            story.append(Paragraph(f"{i}. {safe(r.get('domain',''))}  —  {safe(g.get('한줄요약',''))}", body_st))
            story.append(Paragraph(safe(r.get("url","")[:80]), caption_st))
    doc.build(story)
    return buf.getvalue()


def make_pptx_bytes(all_results, config, research_type, insights, brain_result, grok_signal=""):
    """python-pptx로 한글 PPT 생성"""
    if not _PPT_OK:
        return None
    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    PURPLE = RGBColor(0x63, 0x66, 0xf1)
    DARK   = RGBColor(0x1e, 0x29, 0x3b)
    GRAY   = RGBColor(0x64, 0x74, 0x8b)
    blank_layout = prs.slide_layouts[6]
    def add_slide(title_text, body_lines, title_color=None, small=False):
        if title_color is None: title_color = PURPLE
        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)
        tx = slide.shapes.add_textbox(Emu(457200), Emu(228600), Emu(8229600), Emu(685800))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20 if not small else 16)
        p.font.bold = True
        p.font.color.rgb = title_color
        bx = slide.shapes.add_textbox(Emu(457200), Emu(1066800), Emu(8229600), Emu(3810000))
        bf = bx.text_frame
        bf.word_wrap = True
        for j, line in enumerate(body_lines):
            para = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            para.text = str(line)
            para.font.size = Pt(12 if not small else 10)
            para.font.color.rgb = RGBColor(0xf1, 0xf5, 0xf9)
            para.space_after = Pt(4)
        return slide
    kw = config.get("keyword", "")
    add_slide("리서치 보고서", [f"주제: {kw}", f"타입: {research_type}타입  |  수집: {len(all_results)}개", f"생성: {datetime.now().strftime('%Y-%m-%d %H:%M')}"])
    if insights:
        lines = [l for l in insights.split(chr(10)) if l.strip()]
        chunks, chunk = [], []
        for line in lines:
            chunk.append(line.strip())
            if len(chr(10).join(chunk)) > 500:
                chunks.append(chunk); chunk = []
        if chunk: chunks.append(chunk)
        for ci, ch in enumerate(chunks[:4], 1):
            add_slide(f"종합 인사이트 ({ci}/{min(len(chunks),4)})", ch, small=True)
    if grok_signal:
        glines = [l.strip() for l in grok_signal.split(chr(10)) if l.strip()][:10]
        add_slide("실시간 시장 신호 (Grok)", glines, small=True)
    if brain_result:
        br_lines = []
        j_txt = brain_result.get("judgment","")
        a_txt = brain_result.get("action","")
        if j_txt: br_lines += ["[판단]", j_txt.strip()]
        if a_txt: br_lines += ["", "[액션 플랜]", a_txt.strip()]
        add_slide("뇌 에이전트 마케팅 전략", br_lines, small=True)
    if all_results:
        for ci in range(0, min(len(all_results), 36), 12):
            chunk = all_results[ci:ci+12]
            lines = [f"{ci+j+1}. {r.get('domain','')}  — {r.get('gpt',{}).get('한줄요약','')[:50]}" for j, r in enumerate(chunk)]
            add_slide(f"수집 사이트 ({ci+1}~{ci+len(chunk)})", lines, title_color=GRAY, small=True)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

def make_xlsx_bytes(all_results, config, research_type, insights, brain_result=None):
    needs=config.get("needs",[])
    THIN=Side(style="thin",color="CCCCCC")
    def C(ws,r,c,v,bold=False,bg=None,fc="000000",wrap=True,sz=10,al="left"):
        cell=ws.cell(row=r,column=c,value=v)
        cell.font=Font(bold=bold,color=fc,size=sz,name="맑은 고딕")
        cell.alignment=Alignment(horizontal=al,vertical="top",wrap_text=wrap)
        cell.border=Border(left=THIN,right=THIN,top=THIN,bottom=THIN)
        if bg: cell.fill=PatternFill("solid",fgColor=bg)
        return cell

    wb=openpyxl.Workbook()
    ws0=wb.active; ws0.title="종합 인사이트"
    ws0.merge_cells("A1:B1")
    h=ws0.cell(row=1,column=1,value=f"종합 인사이트 — {config.get('keyword','')}")
    h.font=Font(bold=True,size=14,color="FFFFFF",name="맑은 고딕")
    h.fill=PatternFill("solid",fgColor="1F3864")
    h.alignment=Alignment(horizontal="center",vertical="center")
    ws0.row_dimensions[1].height=30; ws0.column_dimensions["A"].width=120
    for ri,line in enumerate(insights.split("\n"),2):
        c=ws0.cell(row=ri,column=1,value=line)
        c.font=Font(size=11,name="맑은 고딕")
        c.alignment=Alignment(vertical="top",wrap_text=True)
        ws0.row_dimensions[ri].height=18

    if brain_result:
        ws1=wb.create_sheet("뇌 에이전트 전략")
        ws1.merge_cells("A1:B1")
        bh=ws1.cell(row=1,column=1,value=f"마케팅 전략 — {config.get('keyword','')}")
        bh.font=Font(bold=True,size=14,color="FFFFFF",name="맑은 고딕")
        bh.fill=PatternFill("solid",fgColor="2D1B6B")
        bh.alignment=Alignment(horizontal="center",vertical="center")
        ws1.row_dimensions[1].height=30
        ws1.column_dimensions["A"].width=14; ws1.column_dimensions["B"].width=100
        for ri,(lbl,val) in enumerate([
            ("판단",brain_result.get("judgment","")),
            ("이유",brain_result.get("reason","")),
            ("액션",brain_result.get("action","")),
        ],2):
            C(ws1,ri,1,lbl,bold=True,bg="EDE9FE",al="center")
            C(ws1,ri,2,val); ws1.row_dimensions[ri].height=70

    ws=wb.create_sheet("리서치 결과")
    fixed=([("No",5),("페이지명",20),("URL",32),("페이지 제목",26),("한줄 요약",30)]
           if research_type=="A" else
           [("No",5),("출처",8),("URL",32),("도메인",18),("페이지 제목",26),("한줄 요약",30)])
    hdrs=fixed+[(n,30) for n in needs]+[("헤딩",22)]
    nc=len(hdrs)
    ws.merge_cells(start_row=1,start_column=1,end_row=1,end_column=nc)
    t=ws.cell(row=1,column=1,
        value=f"웹 리서치 [{research_type}타입] — {config.get('keyword','')} | {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    t.font=Font(bold=True,size=13,color="FFFFFF",name="맑은 고딕")
    t.fill=PatternFill("solid",fgColor="1F3864")
    t.alignment=Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[1].height=28
    for ci,(h,w) in enumerate(hdrs,1):
        C(ws,2,ci,h,bold=True,bg="2F5496",fc="FFFFFF",al="center")
        ws.column_dimensions[get_column_letter(ci)].width=w
    ws.row_dimensions[2].height=22
    for ri,data in enumerate(all_results,1):
        er=ri+2; bg="EBF0FA" if ri%2 else "FFFFFF"
        g=data.get("gpt",{}); p=data.get("page",{})
        ws.row_dimensions[er].height=80
        vs=([ri,data.get("title",""),data.get("url",""),p.get("page_title",""),g.get("한줄요약","")]
            if research_type=="A" else
            [ri,data.get("source",""),data.get("url",""),data.get("domain",""),p.get("page_title",""),g.get("한줄요약","")])
        for n in needs: vs.append(g.get(n,""))
        vs.append(p.get("headings",""))
        for ci,v in enumerate(vs,1):
            C(ws,er,ci,str(v) if v else "",bg=bg,al="center" if ci==1 else "left")
    ws.freeze_panes="B3"
    buf=io.BytesIO(); wb.save(buf); buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────
#  세션 초기화
# ─────────────────────────────────────────────
for k,v in dict(
    authenticated=False, username="", user_name="", auth_mode="login",
    phase="q1", keyword="", suggested_needs=[], suggested_context="",
    config={}, plan={}, results=None, insights="", research_type="B",
    brain_result=None, fired_patterns=[], feedback_done=False,
).items():
    if k not in st.session_state:
        st.session_state[k] = v

# WebSocket 재연결 시 URL 파라미터로 세션 복구
if not st.session_state.authenticated:
    _p = st.query_params.get("u", "")
    if _p:
        _users = load_users()
        if _p in _users:
            st.session_state.authenticated = True
            st.session_state.username  = _p
            st.session_state.user_name = _users[_p].get("name", _p)


# ═════════════════════════════════════════════
#  인증 화면
# ═════════════════════════════════════════════
if not st.session_state.authenticated:
    st.markdown("""
    <div class="app-hdr" style="text-align:center;">
      <h1>🔍 웹 리서치 어시스턴트</h1>
      <p>GPT 기반 자동 웹 조사 · 분석 · 맞춤형 전략 리포트</p>
    </div>""", unsafe_allow_html=True)

    _, mid, _ = st.columns([1,1.4,1])
    with mid:
        def _set_mode():
            st.session_state.auth_mode = "login" if st.session_state._auth_radio=="로그인" else "register"
        st.radio("로그인 / 계정 만들기", ["로그인","계정 만들기"],
                 index=0 if st.session_state.auth_mode=="login" else 1,
                 horizontal=True, label_visibility="collapsed",
                 key="_auth_radio", on_change=_set_mode)

        with st.container(border=True):
            if st.session_state.auth_mode == "login":
                st.markdown("#### 로그인")
                u = st.text_input("이름", key="li_u", placeholder="가입할 때 입력한 이름")
                p = st.text_input("암호", type="password", key="li_p")
                if st.button("로그인", use_container_width=True):
                    if u and p:
                        user = authenticate(u.strip(), p)
                        if user:
                            st.session_state.authenticated = True
                            st.session_state.username  = u.strip()
                            st.session_state.user_name = user["name"]
                            st.query_params["u"] = u.strip()
                            st.rerun()
                        else:
                            st.error("이름 또는 암호가 맞지 않습니다.")
                    else:
                        st.warning("이름과 암호를 입력하세요.")
            else:
                st.markdown("#### 계정 만들기")
                rn = st.text_input("이름", key="rg_n", placeholder="예: 김철수")
                rp = st.text_input("암호", type="password", key="rg_p")
                rp2= st.text_input("암호 확인", type="password", key="rg_p2")
                if st.button("가입하기", use_container_width=True):
                    if not all([rn, rp, rp2]):
                        st.warning("모든 항목을 입력하세요.")
                    elif rp != rp2:
                        st.error("암호가 일치하지 않습니다.")
                    else:
                        ok, err = register_user(rn, rp)
                        if ok:
                            st.success("계정이 생성됐습니다. 로그인해주세요.")
                            st.session_state.auth_mode = "login"
                            st.rerun()
                        else:
                            st.error(err)
    st.stop()


# ═════════════════════════════════════════════
#  사이드바 — 맞춤 LLM 프로필 표시
# ═════════════════════════════════════════════
with st.sidebar:
    st.markdown(f"### {st.session_state.user_name}")
    if st.button("로그아웃", use_container_width=True):
        st.query_params.clear()
        for k in list(st.session_state.keys()): del st.session_state[k]
        st.rerun()

    # 맞춤 LLM 프로필
    profile = load_profile(st.session_state.username)
    st.markdown("---")
    st.markdown("**🧬 나의 맞춤 LLM 설정**")

    sp = profile.get("system_prompt","")
    if sp:
        st.markdown(f'<div class="profile-box">{sp}</div>', unsafe_allow_html=True)
    else:
        st.caption("아직 피드백이 없습니다.\n리서치 후 피드백을 주시면 다음부터 맞춤형으로 작동합니다.")

    # 적용 중인 규칙 칩
    rules = profile.get("rules", {})
    chips = []
    if rules.get("force_official_only"): chips.append("공식 사이트만")
    if rules.get("exclude_blogs"):       chips.append("블로그 제외")
    if rules.get("min_count"):           chips.append(f"최소 {rules['min_count']}개")
    if rules.get("force_strict_filter"): chips.append(f"필터:{rules['force_strict_filter']}")
    if rules.get("count_multiplier",1)>1:chips.append(f"수량 x{rules['count_multiplier']}")
    if chips:
        chip_html = "".join(f'<span class="rule-chip">{c}</span>' for c in chips)
        st.markdown(chip_html, unsafe_allow_html=True)

    # 최근 조사
    history = load_history(st.session_state.username)
    if history:
        st.markdown("---")
        st.markdown("**최근 조사**")
        for h in history[:5]:
            st.markdown(
                f"<small style='color:#64748b'>{h.get('timestamp','')[:10]}</small><br/>"
                f"<span style='color:#94a3b8;font-size:12px'>{h.get('keyword','')[:18]} ({h.get('result_count',0)}개)</span>",
                unsafe_allow_html=True)

    # 피드백 이력
    fb_history = profile.get("feedback_history",[])
    if fb_history:
        st.markdown("---")
        st.markdown("**최근 피드백**")
        for fb in fb_history[:3]:
            st.markdown(
                f"<small style='color:#64748b'>{fb.get('timestamp','')[:10]}</small><br/>"
                f"<span style='color:#94a3b8;font-size:11px'>{fb.get('changes_summary','')[:40]}</span>",
                unsafe_allow_html=True)


# ═════════════════════════════════════════════
#  헤더
# ═════════════════════════════════════════════
st.markdown(f"""
<div class="app-hdr">
  <h1>🔍 웹 리서치 어시스턴트</h1>
  <p>Claude Sonnet · Opus · Haiku + GPT-4o-mini &nbsp;|&nbsp; {st.session_state.user_name}님 맞춤 설정 적용 중</p>
</div>""", unsafe_allow_html=True)


# ═════════════════════════════════════════════
#  PHASE q1
# ═════════════════════════════════════════════
if st.session_state.phase == "q1":
    with st.container(border=True):
        st.markdown('<div class="step-lbl">Step 1 / 4 — 조사 주제</div>', unsafe_allow_html=True)
        st.markdown("**무엇을 조사할까요?**")
        st.caption("예) 관리감독자 교육 경쟁사  ·  직무교육 협회 목록  ·  AI 자동화 툴 비교")
        kw = st.text_input("kw", value=st.session_state.keyword,
                           placeholder="조사 주제를 입력하세요",
                           label_visibility="collapsed")
    _, mid, _ = st.columns([3,2,3])
    with mid:
        if st.button("분석 시작 →", use_container_width=True):
            if kw.strip():
                st.session_state.keyword = kw.strip()
                _ph = st.empty()
                _ph.info("⏳ AI가 조사 항목을 분석 중입니다...")
                st.session_state.suggested_needs   = gpt_suggest_needs(oai_client, kw.strip())
                st.session_state.suggested_context = gpt_suggest_context(
                    oai_client, kw.strip(), st.session_state.suggested_needs)
                _ph.empty()
                st.session_state.phase = "q_full"
                st.rerun()
            else:
                st.warning("조사 주제를 입력해주세요.")


# ═════════════════════════════════════════════
#  PHASE q_full
# ═════════════════════════════════════════════
elif st.session_state.phase == "q_full":
    with st.container(border=True):
        st.markdown('<div class="step-lbl">Step 1 / 4 ✅</div>', unsafe_allow_html=True)
        st.markdown(f"**{st.session_state.keyword}**")

    with st.container(border=True):
        st.markdown('<div class="step-lbl">Step 2 / 4 — 수집 항목</div>', unsafe_allow_html=True)
        st.markdown("**어떤 정보가 필요한가요?**")
        selected = st.multiselect("항목", options=st.session_state.suggested_needs,
                                  default=st.session_state.suggested_needs,
                                  label_visibility="collapsed")
        custom = st.text_input("직접 추가 항목 (쉼표로 구분)",
                               placeholder="예: 수강 후기, 강사 정보")

    with st.container(border=True):
        st.markdown('<div class="step-lbl">Step 3 / 4 — 추가 조건</div>', unsafe_allow_html=True)
        st.markdown("**추가 조건이나 방향이 있나요?**")
        if st.session_state.suggested_context:
            st.caption(f"💡 AI 제안: {st.session_state.suggested_context}")
        context = st.text_area("조건", value=st.session_state.suggested_context,
                               placeholder="없으면 비워도 됩니다",
                               label_visibility="collapsed", height=80)

    with st.container(border=True):
        st.markdown('<div class="step-lbl">Step 4 / 4 — 수집 개수</div>', unsafe_allow_html=True)
        st.markdown("**몇 개 수집할까요?**")
        st.caption("일반 조사 10~30개  ·  시장 조사 30~100개  ·  대량 수집 100개+")

        # 사용자 규칙에 min_count 있으면 기본값 조정
        rules = load_profile(st.session_state.username).get("rules", {})
        default_count = max(10, rules.get("min_count", 10))
        count = st.number_input("개수", min_value=3, max_value=100,
                                value=default_count, step=5,
                                label_visibility="collapsed")

    st.markdown("")
    c1, c2, c3 = st.columns([2,3,2])
    with c2:
        if st.button("🚀  리서치 시작", use_container_width=True):
            needs = list(selected)
            if custom.strip():
                for item in custom.replace("，",",").split(","):
                    item=item.strip()
                    if item and item not in needs: needs.append(item)
            if not needs: needs = st.session_state.suggested_needs
            st.session_state.config = {
                "keyword": st.session_state.keyword,
                "needs": needs, "context": context.strip(), "count": int(count),
            }
            st.session_state.results      = None
            st.session_state.brain_result = None
            st.session_state.feedback_done= False
            st.session_state.phase        = "running"
            st.rerun()

    if st.button("← 조사 주제 다시 입력"):
        st.session_state.phase = "q1"; st.rerun()


# ═════════════════════════════════════════════
#  PHASE running
# ═════════════════════════════════════════════
elif st.session_state.phase == "running":
    config        = st.session_state.config
    username      = st.session_state.username
    all_results   = []
    insights      = ""
    research_type = "B"
    pages         = []
    driver        = None
    failed        = False
    grok_signal   = ""

    _t0 = time.time()
    def _log(msg): st.write(f"{msg}  ·  {time.time()-_t0:.1f}s")
    with st.status(f"🔍 **{config['keyword']}** 리서치 진행 중...", expanded=True) as status:
        try:
            # 1. 플랜 (DeepSeek deepseek-chat 우선, Claude Sonnet fallback)
            _log("📋 [1/6] 리서치 플랜 수립 중... (DeepSeek R1)")
            plan = build_plan_claude(config, username)
            research_type = plan.get("type","B")
            used_llm = plan.get("_llm","deepseek-chat")
            _log(f"  → {research_type}타입 | {plan.get('reason','')} [{used_llm}]")

            # 2. 사용자 규칙 파이프라인 직접 적용
            rules = load_profile(username).get("rules", {})
            if rules:
                config, plan = apply_user_rules(config, plan, rules)
                applied = []
                if rules.get("force_official_only"): applied.append("공식 사이트 강제")
                if rules.get("min_count"):           applied.append(f"최소 {rules['min_count']}개")
                if rules.get("count_multiplier",1)>1:applied.append(f"수량 x{rules['count_multiplier']}")
                if applied:
                    st.write(f"  🎯 맞춤 규칙 적용: {', '.join(applied)}")

            st.session_state.plan = plan

            # 3. URL 수집
            _log("🌐 [2/6] URL 수집 중... (Naver+Google+DuckDuckGo)")
            if research_type == "A":
                td = plan.get("target_domain","")
                base_url = (td if td.startswith("http") else f"https://www.{td}") if td else ""
                if not base_url:
                    temp = search_all_engines(config, {**plan,"count":1}, oai_client)
                    base_url = temp[0]["url"] if temp else ""
                if not base_url:
                    st.write("  ⚠️ 대상 미발견 → B타입 전환"); research_type="B"
                elif not validate_domain(base_url):
                    st.write("  ⚠️ 접근 불가 → B타입 전환"); research_type="B"
                else:
                    driver = make_driver()
                    home = scrape_page(driver, base_url)
                    if is_login_wall(home):
                        st.write("  🔒 로그인 필요 → B타입 전환"); research_type="B"
                    else:
                        pages = discover_internal_pages(driver, base_url, config, oai_client, config["count"])
                        st.write(f"  ↳ 내부 페이지 {len(pages)}개 선별")

            if research_type in ("B","C"):
                pages = search_all_engines(config, plan, oai_client)
                if not pages: failed = True
                else: st.write(f"  ↳ {len(pages)}개 URL 확정")

            if failed:
                st.error("수집된 URL이 없습니다.")
                status.update(label="URL 수집 실패", state="error")
            else:
                # 4. 스크래핑 + 분석 (GPT-4o-mini)
                _log(f"🤖 [3/6] {len(pages)}개 페이지 분석 중... (GPT-4o-mini)")
                if driver is None: driver = make_driver()
                prog = st.progress(0)

                for i, pg in enumerate(pages):
                    prog.progress((i+1)/len(pages))
                    _log(f"  [{i+1}/{len(pages)}] {pg.get('domain',pg.get('url',''))[:50]}")
                    page_data = scrape_page(driver, pg["url"])
                    if page_data.get("error"):
                        st.write("    ⚠️ 접속 실패 — 건너뜀")
                        try: driver.quit(); driver=make_driver()
                        except: pass
                        continue
                    if is_login_wall(page_data):
                        st.write("    🔒 로그인 장벽 — 건너뜀"); continue
                    gpt_res = analyze_with_gpt(oai_client, pg["url"], page_data, config)
                    st.write(f"    ✅ {gpt_res.get('한줄요약','')[:48]}")
                    all_results.append({**pg, "page": page_data, "gpt": gpt_res})
                    time.sleep(random.uniform(1.0, 1.8))

                # 5. 종합 인사이트 (Claude Opus 4.6 + 개인 시스템 프롬프트)
                if all_results:
                    _log("💡 [4/6] 종합 인사이트 생성 중... (Claude Opus 4.6)")
                    insights = generate_insights(all_results, config, username)

                # 6. Grok 실시간 시장 신호 (xAI grok-3)
                grok_signal = ""
                if all_results and insights:
                    _log("⚡ [5/6] 실시간 시장 신호 보완 중... (Grok)")
                    grok_signal = get_grok_realtime(config, insights)
                    if grok_signal:
                        st.write("  ✅ Grok 실시간 보완 완료")
                    else:
                        st.write("  ⚠️ Grok 연결 실패 (계속 진행)")

                # 7. 뇌 에이전트
                brain_result, fired_patterns = None, []
                if all_results and is_marketing(config):
                    _log("🧠 [6/6] 뇌 에이전트 마케팅 전략 도출 중...")
                    brain_result, fired_patterns = call_brain_agent(config, all_results, insights)
                    if brain_result:
                        st.write(f"  ✅ {len(fired_patterns)}개 패턴 발동")
                    else:
                        st.write("  ⚠️ 뇌 에이전트 연결 실패")

                st.session_state.brain_result   = brain_result
                st.session_state.fired_patterns = fired_patterns
                status.update(label=f"✅ 완료!  {len(all_results)}개 분석", state="complete")

        except Exception as e:
            status.update(label=f"오류: {e}", state="error")
            st.exception(e)
        finally:
            if driver:
                try: driver.quit()
                except: pass

    st.session_state.results       = all_results
    st.session_state.insights      = insights
    st.session_state.grok_signal   = grok_signal
    st.session_state.research_type = research_type
    if all_results:
        save_history(username, {
            "timestamp":    datetime.now().isoformat(),
            "keyword":      config.get("keyword",""),
            "type":         research_type,
            "result_count": len(all_results),
        })
    st.session_state.phase = "q_full" if failed else "results"
    st.rerun()


# ═════════════════════════════════════════════
#  PHASE results
# ═════════════════════════════════════════════
elif st.session_state.phase == "results":
    config        = st.session_state.config
    username      = st.session_state.username
    all_results   = st.session_state.results or []
    insights      = st.session_state.insights
    research_type = st.session_state.research_type
    brain_result  = st.session_state.brain_result
    fired_patterns= st.session_state.fired_patterns or []

    m1,m2,m3,m4 = st.columns(4)
    m1.metric("수집 사이트",  f"{len(all_results)}개")
    m2.metric("조사 타입",    f"{research_type}타입")
    m3.metric("수집 항목",    f"{len(config.get('needs',[]))}개")
    m4.metric("LLM",         "Sonnet+Opus+mini")
    st.markdown("---")

    # 종합 인사이트
    if insights:
        with st.container(border=True):
            st.markdown('<div class="ins-title">💡 Claude Opus 4.6 종합 인사이트</div>', unsafe_allow_html=True)
            st.markdown(insights)

    # 뇌 에이전트
    if brain_result:
        with st.container(border=True):
            st.markdown('<div class="brain-badge">🧠 뇌 에이전트 — 마케팅 전략</div>', unsafe_allow_html=True)
            c1,c2 = st.columns([1,2])
            with c1:
                st.markdown("**판단**"); st.info(brain_result.get("judgment",""))
            with c2:
                st.markdown("**이유**"); st.markdown(brain_result.get("reason",""))
            st.markdown("**액션 아이템**"); st.success(brain_result.get("action",""))
            if fired_patterns:
                with st.expander(f"발동된 패턴 {len(fired_patterns)}개"):
                    for p in fired_patterns:
                        st.markdown(f"- **[{p.get('id','')}]** {p.get('rule','')}")

    # 결과 테이블
    st.markdown("### 📊 수집 결과")
    df = pd.DataFrame()
    if all_results:
        df = results_to_df(all_results, config)
        st.dataframe(df, use_container_width=True, height=360)
        with st.expander("🃏 카드 뷰"):
            for data in all_results:
                g=data.get("gpt",{})
                st.markdown(
                    f'<div class="res-card">'
                    f'<div class="res-domain">{data.get("domain","")}</div>'
                    f'<div class="res-url">{data.get("url","")[:80]}</div>'
                    f'<div class="res-sum">{g.get("한줄요약","")}</div>'
                    f'</div>', unsafe_allow_html=True)

    # 내보내기
    st.markdown("---")
    st.markdown("### 💾 내보내기")
    safe_kw  = re.sub(r'[\\/:*?"<>|]',"_",config.get("keyword","리서치"))
    filename = f"리서치_{safe_kw}_{datetime.now().strftime('%Y%m%d_%H%M')}"
    t1,t2,t3,t4,t5 = st.tabs(["📊 Excel","📄 CSV","{ } JSON","📑 PDF","📌 PPT"])
    with t1:
        st.caption("인사이트 + 뇌 에이전트 전략 + 결과 — 스타일드 Excel")
        if all_results:
            xlsx=make_xlsx_bytes(all_results,config,research_type,insights,brain_result)
            st.download_button("⬇ Excel 다운로드",data=xlsx,
                file_name=f"{filename}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                use_container_width=True)
    with t2:
        if not df.empty:
            csv=df.to_csv(index=False,encoding="utf-8-sig").encode("utf-8-sig")
            st.download_button("⬇ CSV 다운로드",data=csv,
                file_name=f"{filename}.csv",mime="text/csv",use_container_width=True)
    with t3:
        if all_results:
            needs=config.get("needs",[])
            j=json.dumps({"keyword":config.get("keyword"),"needs":needs,"type":research_type,
                "insights":insights,
                "brain_strategy":{"judgment":brain_result.get("judgment","") if brain_result else "",
                                  "action":brain_result.get("action","") if brain_result else ""} if brain_result else None,
                "results":[{"url":r.get("url",""),"domain":r.get("domain",""),
                    "summary":r.get("gpt",{}).get("한줄요약",""),
                    **{n:r.get("gpt",{}).get(n,"") for n in needs}} for r in all_results]
            },ensure_ascii=False,indent=2).encode("utf-8")
            st.download_button("⬇ JSON 다운로드",data=j,
                file_name=f"{filename}.json",mime="application/json",use_container_width=True)

    with t4:
        if not _PDF_OK:
            st.warning("PDF 라이브러리 미설치 (pip install reportlab)")
        else:
            grok_sig = st.session_state.get("grok_signal","")
            st.caption("한글 PDF — 인사이트 + Grok 신호 + 뇌 에이전트 + 사이트 목록")
            if st.button("PDF 생성 및 다운로드", use_container_width=True, key="pdf_btn"):
                with st.spinner("PDF 생성 중..."):
                    pdf_bytes = make_pdf_bytes(all_results, config, research_type, insights, brain_result, grok_sig)
                if pdf_bytes:
                    st.download_button("⬇ PDF 저장", data=pdf_bytes,
                        file_name=f"{filename}.pdf", mime="application/pdf",
                        use_container_width=True)
                else:
                    st.error("PDF 생성 실패")
    with t5:
        if not _PPT_OK:
            st.warning("PPT 라이브러리 미설치 (pip install python-pptx)")
        else:
            grok_sig = st.session_state.get("grok_signal","")
            st.caption("16:9 슬라이드 — 표지 + 인사이트 + 뇌 에이전트 + 사이트 목록")
            if st.button("PPT 생성 및 다운로드", use_container_width=True, key="ppt_btn"):
                with st.spinner("PPT 생성 중..."):
                    pptx_bytes = make_pptx_bytes(all_results, config, research_type, insights, brain_result, grok_sig)
                if pptx_bytes:
                    st.download_button("⬇ PPT 저장", data=pptx_bytes,
                        file_name=f"{filename}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True)
                else:
                    st.error("PPT 생성 실패")

    # ── 진짜 맞춤형 피드백 ──────────────────────────
    st.markdown("---")
    st.markdown("### 💬 피드백 — 다음 리서치에 즉각 반영됩니다")
    st.caption("불만·요청 사항을 자유롭게 말해주세요. Claude Haiku가 즉시 분석해서 당신만의 LLM 설정을 업데이트합니다.")

    if st.session_state.feedback_done:
        st.success("✅ 피드백이 반영됐습니다. 다음 리서치부터 맞춤 설정이 적용됩니다.")
        updated_profile = load_profile(username)
        if updated_profile.get("system_prompt"):
            st.markdown("**업데이트된 나의 맞춤 LLM 설정:**")
            st.markdown(f'<div class="profile-box">{updated_profile["system_prompt"]}</div>',
                        unsafe_allow_html=True)
    else:
        fb_text = st.text_area(
            "피드백 입력",
            placeholder=(
                "예시:\n"
                "• 결과가 너무 적었어, 항상 30개 이상 줘\n"
                "• 블로그 말고 공식 사이트만 줘\n"
                "• 쿼리가 너무 좁았어, 더 다양하게 쳐줘\n"
                "• 인사이트가 너무 일반적이야, 더 구체적으로\n"
                "• 다음부터 뉴스도 포함해줘"
            ),
            height=140,
            label_visibility="collapsed",
        )
        if st.button("🔄 반영하기", use_container_width=True):
            if fb_text.strip():
                with st.spinner("Claude Haiku가 피드백을 분석하고 당신의 LLM을 업데이트하는 중..."):
                    result = process_feedback(username, fb_text.strip())
                st.session_state.feedback_done = True

                # 변경 사항 표시
                changes = result.get("changes", [])
                if changes:
                    st.markdown("**적용된 변경 사항:**")
                    for ch in changes:
                        st.markdown(f'<div class="change-item">✓ {ch}</div>',
                                    unsafe_allow_html=True)
                st.rerun()
            else:
                st.warning("피드백 내용을 입력해주세요.")

    # 액션 버튼
    st.markdown("---")
    a1, a2 = st.columns(2)
    with a1:
        if st.button("🔄 새 리서치 시작", use_container_width=True):
            for k in ["phase","keyword","suggested_needs","suggested_context",
                      "config","plan","results","insights","research_type",
                      "brain_result","fired_patterns","feedback_done"]:
                st.session_state.pop(k, None)
            st.session_state.phase = "q1"
            st.rerun()
    with a2:
        if st.button("⚙️ 같은 주제로 재조사", use_container_width=True):
            st.session_state.results       = None
            st.session_state.brain_result  = None
            st.session_state.feedback_done = False
            st.session_state.phase         = "q_full"
            st.rerun()
